#!/usr/bin/env python3
"""Extract raw corporate / research source files into readable markdown.

Part of the `brooker-db-to-wiki` workflow. Walks a source root, converts each
docx/pdf/xlsx/pptx/doc file to plain markdown, and writes the result to a
gitignored scratch dir so wiki articles can be authored from real content.

Each extract records the source file's last-modified timestamp, so the wiki
articles can be dated for later recall.

Source roots are READ ONLY — this script never writes to them.

Usage:
    python scripts/extract_sources.py ceo
    python scripts/extract_sources.py all
    python scripts/extract_sources.py --root "O:/2nd_Brain" Crypto
    python scripts/extract_sources.py --root "O:/2nd_Brain" --out .source-extracts/2nd_brain all

Requires: pdfplumber, python-docx, openpyxl, python-pptx, xlrd
(.doc files fall back to the `antiword` CLI if available.)
"""
from __future__ import annotations

import argparse
import datetime as _dt
import shutil
import subprocess
import sys
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parent.parent
DEFAULT_ROOT = Path("O:/brooker_database")
SKIP_EXT = {".jpg", ".jpeg", ".png", ".gif", ".db", ".ini", ".php", ".css", ".js"}


def extract_pdf(path: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            parts.append(f"\n\n--- page {i} ---\n\n{text}")
            for table in page.extract_tables():
                rows = ["| " + " | ".join((c or "").strip() for c in row) + " |" for row in table]
                if rows:
                    parts.append("\n" + "\n".join(rows))
    return "".join(parts)


def extract_docx(path: Path) -> str:
    import docx

    document = docx.Document(str(path))
    parts: list[str] = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            parts.append("| " + " | ".join(cells) + " |")
    return "\n\n".join(parts)


def extract_doc(path: Path) -> str:
    """Old-format .doc via the antiword CLI (if installed)."""
    if shutil.which("antiword") is None:
        raise RuntimeError("antiword not available for legacy .doc")
    out = subprocess.run(
        ["antiword", str(path)], capture_output=True, text=True, encoding="utf-8", errors="replace"
    )
    if out.returncode != 0:
        raise RuntimeError(f"antiword failed: {out.stderr.strip()}")
    return out.stdout


def extract_xlsx(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"\n\n## sheet: {sheet.title}\n")
        for row in sheet.iter_rows(values_only=True):
            if any(v is not None for v in row):
                cells = ["" if v is None else str(v) for v in row]
                parts.append("| " + " | ".join(cells) + " |")
    wb.close()
    return "\n".join(parts)


def extract_xls(path: Path) -> str:
    import xlrd

    book = xlrd.open_workbook(str(path))
    parts: list[str] = []
    for sheet in book.sheets():
        parts.append(f"\n\n## sheet: {sheet.name}\n")
        for r in range(sheet.nrows):
            cells = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
            if any(c.strip() for c in cells):
                parts.append("| " + " | ".join(cells) + " |")
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"\n\n--- slide {i} ---\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    parts.append("| " + " | ".join(cells) + " |")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"\n[speaker notes]\n{notes}")
    return "\n\n".join(parts)


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".doc": extract_doc,
    ".xlsx": extract_xlsx,
    ".xls": extract_xls,
    ".pptx": extract_pptx,
}


def extract_folder(root: Path, sub: str, out_root: Path) -> None:
    src_dir = root / sub
    if not src_dir.is_dir():
        print(f"  [skip] {sub}: no source folder at {src_dir}")
        return

    out_dir = out_root / sub
    out_dir.mkdir(parents=True, exist_ok=True)

    for path in sorted(p for p in src_dir.rglob("*") if p.is_file()):
        ext = path.suffix.lower()
        rel = path.relative_to(src_dir)
        if ext in SKIP_EXT:
            continue
        out_path = out_dir / rel.with_suffix(rel.suffix + ".md")
        out_path.parent.mkdir(parents=True, exist_ok=True)
        mtime = _dt.date.fromtimestamp(path.stat().st_mtime).isoformat()
        header = f"# {path.name}\n\n_Source: `{path}`_\n_Modified: {mtime}_\n\n"
        extractor = EXTRACTORS.get(ext)
        if extractor is None:
            out_path.write_text(
                header + f"> Unsupported format `{ext}` — extract manually.\n",
                encoding="utf-8",
            )
            print(f"  [manual] {rel} — unsupported {ext}")
            continue
        try:
            content = extractor(path)
        except Exception as exc:  # noqa: BLE001 - report and continue
            out_path.write_text(header + f"> Extraction failed: {exc}\n", encoding="utf-8")
            print(f"  [error] {rel} — {exc}")
            continue
        out_path.write_text(header + content + "\n", encoding="utf-8")
        print(f"  [ok]    {rel}  (modified {mtime})")


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(description="Extract source files to markdown.")
    parser.add_argument("target", help="subfolder name under the root, or 'all'")
    parser.add_argument("--root", default=str(DEFAULT_ROOT), help="source root directory")
    parser.add_argument("--out", default=None, help="output dir (default .source-extracts/<root-name>)")
    args = parser.parse_args(argv[1:])

    root = Path(args.root)
    if not root.is_dir():
        print(f"Source root not found: {root}")
        return 1
    out_root = Path(args.out) if args.out else REPO_ROOT / ".source-extracts" / root.name.lower()

    subs = (
        sorted(p.name for p in root.iterdir() if p.is_dir())
        if args.target == "all"
        else [args.target]
    )
    for sub in subs:
        print(f"Extracting {root.name}/{sub}/ ...")
        extract_folder(root, sub, out_root)
    print(f"\nDone. Extracts written to {out_root}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
