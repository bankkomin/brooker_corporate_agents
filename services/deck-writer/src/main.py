"""deck-writer — generate a PowerPoint pitch deck from a brief.

Flow:
  POST /compose
    body: { "brief": "...", "dept_id": "ic", "title": "...", "audience": "..." }
    1. Embed brief via rag-ingestion /embed
    2. Search Qdrant collections relevant to dept_id
    3. LLM (Qwen) composes a slide JSON spec from brief + context
    4. python-pptx renders the spec into a .pptx (template if available)
    5. Save to /data/decks/<id>.pptx, return URL/path

  POST /report
    Produces a .docx review document from the same drafter pipeline.
    Optional rich-content fields: images, charts, mermaid diagrams.

  POST /compose-xlsx
    Builds a .xlsx workbook from an XlsxComposeSpec.
"""
from __future__ import annotations

import asyncio
import json
import logging
import os
import re
import tempfile
import time
import uuid
from contextlib import asynccontextmanager
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Literal

import httpx
import structlog
from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from pydantic import BaseModel, Field

# ----- rich-content helper modules (frozen, tested) -----
from services.shared.chart_render import render_chart_png, validate_chart_spec
from services.shared.drafter_table_prompt import (
    extract_tables_from_text,
    table_emission_prompt_snippet,
)
from services.shared.image_embed import (
    ImageSource,
    embed_image_in_docx,
    embed_image_in_pptx,
)
from services.shared.mermaid_render import render_mermaid_png
from services.shared.table_render import (
    add_table_to_docx,
    add_table_to_pptx,
    validate_table_spec,
)
from services.shared.xlsx_compose import compose_xlsx, validate_xlsx_spec

load_dotenv()
logging.basicConfig(level=os.getenv("LOG_LEVEL", "INFO").upper())
log = structlog.get_logger("deck-writer")

# ----- config -----
LLM_BASE_URL = os.getenv("VLLM_LARGE_URL", "http://nginx:8080/v1").rstrip("/")
LLM_MODEL = os.getenv("VLLM_LARGE_MODEL", "qwen-large")
LLM_API_KEY = os.getenv("LLM_API_KEY", "")
# The DGX Spark fails beyond a few concurrent sequences — cap concurrent LLM calls.
LLM_MAX_CONCURRENCY = int(os.getenv("LLM_MAX_CONCURRENCY", "4"))
_DGX_SEMAPHORE = asyncio.Semaphore(LLM_MAX_CONCURRENCY)
RAG_INGESTION_URL = os.getenv("RAG_INGESTION_URL", "http://rag-ingestion:3004").rstrip("/")
QDRANT_URL = os.getenv("QDRANT_URL", "http://qdrant:6333").rstrip("/")
DEPARTMENTS_CONFIG = os.getenv("DEPARTMENTS_CONFIG", "/app/config/templates/../departments.json")

# Per-dept .pptx template. Falls back to DECK_TEMPLATE_PATH (legacy IC default)
# when the dept isn't in the map, so existing callers keep working.
_TEMPLATES_ROOT = "/app/config/templates"
_TEMPLATE_BY_DEPT: dict[str, str] = {
    "ic":  f"{_TEMPLATES_ROOT}/ic/IC-meeting-deck-reference.pptx",
    "vcc": f"{_TEMPLATES_ROOT}/vcc/vcc-deck-template.pptx",
    "cio": f"{_TEMPLATES_ROOT}/vcc/vcc-deck-template.pptx",   # CIO shares VCC look
}
_BROOKER_TEMPLATE = f"{_TEMPLATES_ROOT}/brooker/brooker-deck-template.pptx"

TEMPLATE_PATH = os.getenv("DECK_TEMPLATE_PATH", _BROOKER_TEMPLATE)


def _template_for(dept_id: str) -> str:
    """Resolve which .pptx template to load. Honor env override first."""
    env_override = os.getenv("DECK_TEMPLATE_PATH")
    if env_override:
        return env_override
    return _TEMPLATE_BY_DEPT.get(dept_id, _BROOKER_TEMPLATE)
OUTPUT_DIR = Path(os.getenv("DECK_OUTPUT_DIR", "/data/decks"))
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
RAG_TOP_K = int(os.getenv("RAG_TOP_K", "8"))
RAG_MIN_RELEVANCE = float(os.getenv("RAG_MIN_RELEVANCE", "0.50"))


# ----- rich-content embedding models (shared by /compose and /report) -----

class ImageEmbed(BaseModel):
    """Embed a PNG/JPEG image by URL, filesystem path, or MinIO key."""
    url: str | None = None
    path: str | None = None
    minio_key: str | None = None
    section_hint: str | None = None        # docx: heading-substring match
    slide_index: int | None = None         # pptx: 0-based slide index
    slide_title_hint: str | None = None    # pptx: slide title substring match
    caption: str | None = None
    width_inches: float = 5.5


class ChartEmbed(BaseModel):
    """Render a matplotlib chart from a ChartSpec dict and embed as PNG."""
    spec: dict                             # validated by validate_chart_spec
    section_hint: str | None = None
    slide_index: int | None = None
    slide_title_hint: str | None = None
    caption: str | None = None
    width_inches: float = 5.5


class MermaidEmbed(BaseModel):
    """Render a Mermaid diagram and embed as PNG."""
    mermaid: str
    theme: str = "default"
    background: str = "white"
    section_hint: str | None = None
    slide_index: int | None = None
    slide_title_hint: str | None = None
    caption: str | None = None
    width_inches: float = 5.5


# ----- request / response models -----
class ComposeRequest(BaseModel):
    # Accept either `brief` (preferred) or `query` (slack-bot router shape).
    brief: str | None = Field(default=None, description="What the deck is about (1-3 sentences).")
    dept_id: str = "ic"
    title: str | None = None
    audience: str | None = None
    user_id: str | None = None
    # match the shape used by other orchestrators so slack-bot can route here:
    query: str | None = None
    channel: str | None = None
    thread_ts: str | None = None
    # rich-content embeds — all optional; default to empty list for back-compat
    images: list[ImageEmbed] = Field(default_factory=list)
    charts: list[ChartEmbed] = Field(default_factory=list)
    mermaid: list[MermaidEmbed] = Field(default_factory=list)


class ComparisonColumn(BaseModel):
    heading: str = ""
    bullets: list[str] = Field(default_factory=list)


class KpiStat(BaseModel):
    value: str            # the big number / metric, e.g. "100%", "Bt 1.2bn", ">30 days"
    label: str            # short caption under the value
    note: str | None = None  # optional sub-note (source / caveat)


class SlideSpec(BaseModel):
    # layout drives the renderer; defaults to plain bullets for back-compat.
    layout: Literal["cover", "section", "bullets", "comparison", "kpi_stats",
                    "quote", "closing"] = "bullets"
    title: str = ""
    bullets: list[str] = Field(default_factory=list)
    notes: str | None = None
    # layout-specific extras (only used when layout matches)
    subtitle: str | None = None              # cover / closing
    columns: list[ComparisonColumn] = Field(default_factory=list)  # comparison
    stats: list[KpiStat] = Field(default_factory=list)             # kpi_stats
    quote: str | None = None                 # quote
    attribution: str | None = None           # quote


class DeckSpec(BaseModel):
    title: str = "Untitled Deck"
    subtitle: str | None = None
    # Default empty list so Qwen omitting the field doesn't kill the request;
    # the renderer adds a cover and we fall through to a single "Context limited"
    # slide if no slides came back.
    slides: list[SlideSpec] = Field(default_factory=list)


class ComposeResponse(BaseModel):
    answer: str  # slack-bot shape: textual summary the user sees
    confidence: str = "Medium"
    error: str | None = None
    file_path: str  # absolute path inside container
    file_name: str  # basename for nice display
    file_url: str | None = None  # if exposed via static route
    sources: list[dict] = Field(default_factory=list)


class ComposeXlsxRequest(BaseModel):
    """Request body for POST /compose-xlsx."""
    spec: dict = Field(description="XlsxComposeSpec — sheets, optional charts, optional metadata.")
    filename: str = Field(default="report.xlsx", description="Output filename (basename only).")
    caller_dept: str = Field(default="", description="Requesting department; used for output path scoping.")


# ----- HTTP clients -----
_http: httpx.AsyncClient | None = None


@asynccontextmanager
async def lifespan(app: FastAPI):
    global _http
    headers = {"Content-Type": "application/json"}
    if LLM_API_KEY:
        headers["Authorization"] = f"Bearer {LLM_API_KEY}"
    _http = httpx.AsyncClient(timeout=120.0, headers=headers)
    log.info("deck-writer.startup", port=3050, output_dir=str(OUTPUT_DIR),
             template_exists=Path(TEMPLATE_PATH).is_file())
    yield
    await _http.aclose()


app = FastAPI(title="deck-writer", version="0.1.0", lifespan=lifespan)


@app.get("/health")
async def health() -> dict[str, Any]:
    return {
        "status": "ok",
        "service": "deck-writer",
        "template_exists": Path(TEMPLATE_PATH).is_file(),
        "output_dir": str(OUTPUT_DIR),
    }


@app.get("/files/{name}")
async def get_file(name: str) -> FileResponse:
    # Whitelist: only filenames we generated (no path traversal).
    safe = re.sub(r"[^A-Za-z0-9_.\-]", "", name)
    p = OUTPUT_DIR / safe
    if not p.is_file():
        raise HTTPException(404, "deck not found")
    return FileResponse(
        path=str(p),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=safe,
    )


# ----- RAG -----
async def _embed(text: str) -> list[float]:
    assert _http is not None
    r = await _http.post(f"{RAG_INGESTION_URL}/embed", json={"text": text})
    r.raise_for_status()
    return r.json()["vector"]


def _resolve_collections(dept_id: str) -> list[str]:
    cfg_path = Path(DEPARTMENTS_CONFIG)
    cross: list[str] = []
    if cfg_path.is_file():
        try:
            data = json.loads(cfg_path.read_text(encoding="utf-8"))
            depts = data.get("departments", {}) or {}
            entry = depts.get(dept_id, {}) or {}
            cross = entry.get("crossReadAccess", []) or []
            if "*" in cross:
                cross = [k for k, v in depts.items() if k != dept_id and v.get("live", False)]
        except Exception as exc:
            log.warning("departments_load_failed", err=str(exc))
    own = [f"{dept_id}_docs", f"{dept_id}_chat", f"{dept_id}_knowledge"]
    return own + ["shared_policies"] + [f"{d}_docs" for d in cross]


async def _qdrant_search(collection: str, vec: list[float]) -> list[dict]:
    assert _http is not None
    try:
        r = await _http.post(
            f"{QDRANT_URL}/collections/{collection}/points/search",
            json={
                "vector": vec, "limit": RAG_TOP_K, "with_payload": True,
                "score_threshold": RAG_MIN_RELEVANCE,
            },
        )
        if r.status_code != 200:
            return []
        return r.json().get("result", []) or []
    except Exception:
        return []


async def _retrieve(brief: str, dept_id: str) -> list[dict]:
    vec = await _embed(brief)
    collections = _resolve_collections(dept_id)
    hits: list[dict] = []
    for coll in collections:
        for p in await _qdrant_search(coll, vec):
            payload = p.get("payload", {}) or {}
            hits.append({
                "source": payload.get("original_filename") or payload.get("source_file") or coll,
                "excerpt": (payload.get("text", "") or "")[:600],
                "score": p.get("score", 0.0),
                "collection": coll,
            })
    hits.sort(key=lambda x: x["score"], reverse=True)
    return hits[: RAG_TOP_K * 2]


# ----- LLM pipeline: two-agent composition -----
#
# Agent 1 (DRAFTER): decides WHAT each slide should communicate.
# Agent 2 (DESIGNER): decides HOW each slide should look (layout assignment).
#
# Rationale: separation of concerns. The drafter focuses on narrative + grounding
# in retrieved context; the designer focuses on visual variety. Each call is
# scoped + smaller, so JSON output is more reliable.

_DRAFTER_SYSTEM = (
    "You are a deck DRAFTER. You decide what each slide must communicate — its "
    "headline, the key points it makes, and the concrete data points it cites. "
    "You do NOT decide layout / format / colour — a separate Designer step "
    "handles that.\n"
    "\n"
    "Output ONLY valid JSON, no prose, no markdown fences:\n"
    "{\n"
    '  "title": "deck-level title",\n'
    '  "subtitle": "string|null",\n'
    '  "outline": [\n'
    '    {\n'
    '      "narrative_role": "cover|context|analysis|deep_dive|comparison|metrics|quote|conclusion",\n'
    '      "headline": "what this slide must communicate in one sentence",\n'
    '      "key_points": ["point 1", "point 2", ...],   // 2-6 short statements\n'
    '      "supporting_data": ["concrete numbers / quotes / citations from context"],\n'
    '      "compare_a": "if comparison: first thing being compared",\n'
    '      "compare_b": "if comparison: second thing",\n'
    '      "quotable":  "if quote: the exact line worth pulling out",\n'
    '      "quotable_source": "if quote: the filename / speaker"\n'
    "    }\n"
    "  ]\n"
    "}\n"
    "\n"
    "Composition rules (READ CAREFULLY — variety matters):\n"
    "- **8-12 outline items total.** First item narrative_role MUST be \"cover\". Last MUST be \"conclusion\".\n"
    "- Ground EVERY key_point + supporting_data in the provided context. Cite source "
    "filenames inline like [filename.pdf]. Never invent figures or quotes.\n"
    "- **Use a MIX of narrative_role values across the outline. Never use only one role.** Aim for:\n"
    "    • 1 \"cover\" (slide 1)\n"
    "    • 1 \"conclusion\" (last slide)\n"
    "    • 2-4 \"analysis\" or \"deep_dive\" (the workhorses)\n"
    "    • At least 1 \"metrics\" — when sources contain 2+ concrete numbers (percentages, ratios, $/Bt amounts, thresholds, dates, counts). Put those numbers in supporting_data so Designer can render them as stat cards.\n"
    "    • At least 1 \"comparison\" — whenever the topic naturally contrasts two things (LCR vs NSFR, Level 1 vs Level 2 HQLA, before vs after, dept A vs dept B, two vendors, two scenarios). Set compare_a + compare_b.\n"
    "    • 0-1 \"quote\" — only if a source contains an attributable, quotable line; set quotable + quotable_source.\n"
    "- Specifically: if you're tempted to make EVERY slide \"analysis\" or \"deep_dive\", you're under-using metrics and comparison. Re-read the context for numbers and contrasts you missed.\n"
    "- If context is thin for a slide, write {\"narrative_role\":\"context\","
    "\"headline\":\"Context limited: <topic>\",\"key_points\":[...]} and explain "
    "what would be needed.\n"
    "\n"
    "GOOD EXAMPLES of narrative_role choice:\n"
    "  Topic: LCR with HQLA tiers\n"
    "    ✓ \"metrics\": supporting_data=[\"LCR minimum: 100% [lcr.md]\", \"Level 1 haircut: 0% [lcr.md]\", \"Level 2A haircut: 15% [lcr.md]\", \"Stress window: 30 days [lcr.md]\"]\n"
    "    ✓ \"comparison\": compare_a=\"Level 1 HQLA\", compare_b=\"Level 2 HQLA\"\n"
    "  Topic: Q1 capital review\n"
    "    ✓ \"metrics\": supporting_data=[\"CAR: 14.2% [tracker.xlsx]\", \"CET1: 11.8% [tracker.xlsx]\", \"RWA: Bt 8.4bn [tracker.xlsx]\"]\n"
    "    ✓ \"comparison\": compare_a=\"Q4 2025\", compare_b=\"Q1 2026\"\n"
    "\n"
    "HEADLINE LENGTH RULES (critical for layout — long headlines overflow into body):\n"
    "- `headline` must be ≤8 words and ≤60 characters. Slide TITLE not a sentence.\n"
    "- ❌ BAD: \"Hex Trust represents a shift from unregulated trading structures to regulated custody.\"\n"
    "- ✓ GOOD: \"From Unregulated OTC to Regulated Custody\"\n"
    "- ❌ BAD: \"The previous Fireblocks arrangement carried significantly higher risk profiles.\"\n"
    "- ✓ GOOD: \"Fireblocks Risk Profile (Prior State)\"\n"
    "- Put the full sentence in `key_points` if it's an argument; the headline is just the section label.\n"
)

# Append the table-emission instructions to the drafter system prompt so that
# the LLM can emit ```table JSON blocks when comparison/grid data would help.
# extract_tables_from_text() parses these blocks out before rendering.
_DRAFTER_SYSTEM = _DRAFTER_SYSTEM + "\n\n" + table_emission_prompt_snippet()

_DESIGNER_SYSTEM = (
    "You are a deck DESIGNER. You receive a content outline and produce the final "
    "JSON spec that the renderer uses. You decide the visual layout for each "
    "slide. You may merge or split adjacent outline items if doing so makes the "
    "deck stronger (e.g. two comparison-ish items → one comparison slide).\n"
    "\n"
    "Available layouts: cover, section, bullets, comparison, kpi_stats, quote, closing.\n"
    "\n"
    "Layout-selection rules:\n"
    "- narrative_role=\"cover\"      → layout=\"cover\"\n"
    "- narrative_role=\"conclusion\" → layout=\"closing\"\n"
    "- narrative_role=\"metrics\" AND supporting_data has 2-4 number-shaped entries → layout=\"kpi_stats\"\n"
    "- narrative_role=\"comparison\" AND compare_a + compare_b present → layout=\"comparison\"\n"
    "- supporting_data contains one quotable sentence with attribution → layout=\"quote\"\n"
    "- otherwise → layout=\"bullets\"\n"
    "- Insert ONE layout=\"section\" divider between major narrative phase changes (every 3-4 slides). Sections you insert must use a short title; key_points/supporting_data can be ignored.\n"
    "\n"
    "Output ONLY valid JSON, no prose, no fences:\n"
    "{\n"
    '  "title": "from outline",\n'
    '  "subtitle": "from outline",\n'
    '  "slides": [{\n'
    '    "layout": "cover|section|bullets|comparison|kpi_stats|quote|closing",\n'
    '    "title": "string",\n'
    '    "subtitle": "string|null",          // cover, closing\n'
    '    "bullets": ["string", ...],          // bullets / section\n'
    '    "columns": [{"heading":"string","bullets":["string",...]}, ...],   // comparison: exactly 2\n'
    '    "stats":   [{"value":"100%","label":"LCR floor","note":"[lcr.md]"}, ...],  // kpi_stats: 2-4\n'
    '    "quote":   "string",                 // quote\n'
    '    "attribution": "string",             // quote\n'
    '    "notes":   "string|null"             // speaker notes\n'
    "  }]\n"
    "}\n"
    "\n"
    "Content-conversion rules:\n"
    "- For bullets/section: bullets come straight from outline.key_points. Each ≤20 words.\n"
    "- For kpi_stats: parse each supporting_data line for a number; value is short (≤15 chars), label is 2-5 words, note cites the source.\n"
    "- For comparison: use compare_a/compare_b as column headings; split key_points + supporting_data evenly across the two columns.\n"
    "- For quote: pick the most quotable supporting_data line; attribution is the cited filename.\n"
    "- Preserve all source filename citations.\n"
)


async def _llm_call(system: str, user: str, *, max_tokens: int = 3000) -> dict | None:
    assert _http is not None
    body = {
        "model": LLM_MODEL,
        "messages": [
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
        "max_tokens": max_tokens,
        "temperature": 0.2,
        "response_format": {"type": "json_object"},
        "chat_template_kwargs": {"enable_thinking": False},
    }
    async with _DGX_SEMAPHORE:
        r = await _http.post(f"{LLM_BASE_URL}/chat/completions", json=body)
        r.raise_for_status()
        raw = r.json()["choices"][0]["message"]["content"] or "{}"
    return _extract_json_object(raw)


async def _draft_outline(brief: str, context: str, title: str | None,
                          audience: str | None) -> dict:
    """Agent 1 — Drafter. Returns a content outline (no layouts)."""
    parts: list[str] = []
    if title:
        parts.append(f"Requested title: {title}")
    if audience:
        parts.append(f"Audience: {audience}")
    parts.append(f"Brief: {brief}")
    parts.append(f"Retrieved context:\n{context or '(no relevant context retrieved)'}")
    t0 = time.monotonic()
    parsed = await _llm_call(_DRAFTER_SYSTEM, "\n\n".join(parts), max_tokens=3000)
    log.info("draft.done", elapsed_ms=int((time.monotonic() - t0) * 1000),
             outline_items=len((parsed or {}).get("outline", []) or []))
    if not parsed:
        raise HTTPException(502, "Drafter returned invalid JSON")
    parsed.setdefault("title", title or "Untitled Deck")
    parsed.setdefault("outline", [])
    return parsed


async def _design_layouts(outline: dict) -> DeckSpec:
    """Agent 2 — Designer. Maps outline → DeckSpec with layouts assigned.

    If the LLM call comes back empty / unparseable / with zero slides, fall
    through to the deterministic mapper so the user still gets a deck.
    """
    user_payload = json.dumps(outline, ensure_ascii=False)
    t0 = time.monotonic()
    parsed = await _llm_call(_DESIGNER_SYSTEM,
                              f"Outline to design:\n{user_payload}",
                              max_tokens=3500)
    parsed_slides = (parsed or {}).get("slides") or []
    log.info("design.done", elapsed_ms=int((time.monotonic() - t0) * 1000),
             slides=len(parsed_slides))

    if not parsed or not parsed_slides:
        # Drafter wrote a real outline but Designer didn't produce slides —
        # don't waste the upstream work, run the deterministic mapper.
        log.warning("design.fallback",
                    reason="empty_slides" if parsed else "designer_returned_invalid_json")
        return _deterministic_design(outline)

    try:
        return DeckSpec.model_validate(parsed)
    except Exception as exc:
        log.warning("design.spec_validation_failed", error=str(exc),
                    fallback="deterministic_design")
        return _deterministic_design(outline)


async def _draft_outline_with_retry(brief: str, context: str, title: str | None,
                                    audience: str | None) -> dict:
    """Draft with one retry if the first call returns an empty outline.

    HTTPException is caught explicitly because in some FastAPI/Starlette builds it
    does not inherit from Exception, so a bare `except Exception` would miss it and
    let the 502 propagate past the minimal-outline fallback.
    """
    try:
        outline = await _draft_outline(brief, context, title, audience)
    except (Exception, HTTPException) as exc:
        log.warning("draft.first_call_failed", error=str(exc))
        outline = {}
    items = outline.get("outline") or []
    if items:
        return outline
    log.warning("draft.empty_outline_retry")
    # Stronger nudge on the retry — explicit "you must include >=4 outline items".
    nudge = (
        "Your previous reply had no outline items. You MUST return 4-8 outline "
        "items in `outline`. Use the schema in the system prompt."
    )
    try:
        outline2 = await _draft_outline(brief, context + "\n\n" + nudge, title, audience)
    except (Exception, HTTPException) as exc:
        log.warning("draft.retry_call_failed", error=str(exc))
        outline2 = {}
    items2 = outline2.get("outline") or []
    if items2:
        return outline2
    # Last resort: synthesise a tiny outline from the brief so the user still gets
    # *something*. The Designer fallback will fill in layout decisions.
    log.warning("draft.synthesised_minimal_outline")
    return {
        "title": title or "Briefing",
        "subtitle": audience,
        "outline": [
            {"narrative_role": "cover", "headline": title or "Briefing", "key_points": [], "supporting_data": []},
            {"narrative_role": "context", "headline": "Brief", "key_points": [brief], "supporting_data": []},
            {"narrative_role": "context", "headline": "Context limited",
             "key_points": [
                 "The model couldn't draft a content outline from the retrieved sources.",
                 "Try a more specific brief or share a source document directly.",
             ], "supporting_data": []},
            {"narrative_role": "conclusion", "headline": "Q&A", "key_points": [], "supporting_data": []},
        ],
    }


_NUM_RE = re.compile(
    r"([+\-]?\$?\€?\£?\¥?\฿?[\d][\d,]*(?:\.\d+)?\s*"
    r"(?:%|bps|bn|mn?|k|x|days?|years?|months?|hours?|wk|qty|t)?)",
    re.IGNORECASE,
)


def _extract_numbers(text: str) -> list[str]:
    """Return numeric-looking tokens in `text`. Filters tiny matches like '1' alone."""
    raw = _NUM_RE.findall(text or "")
    out: list[str] = []
    for t in raw:
        t = t.strip()
        if not t or len(t) < 2:
            continue
        # Skip pure-noise matches like ',', '.'
        if not any(c.isdigit() for c in t):
            continue
        out.append(t)
    return out


def _looks_like_quote(text: str) -> str | None:
    """If `text` contains a clearly quotable sentence, return it; else None."""
    if not text:
        return None
    m = re.search(r'["“](.{40,200}?)["”]', text)
    if m:
        return m.group(1).strip()
    return None


def _strip_citations(text: str) -> str:
    """Remove inline [filename.ext] citations for clean labels."""
    return re.sub(r"\s*\[[^\]]+\]\s*", " ", text or "").strip()


def _first_citation(text: str) -> str | None:
    m = re.search(r"\[([^\]]+)\]", text or "")
    return m.group(1) if m else None


def _kpi_from_lines(lines: list[str], max_n: int = 4) -> list[KpiStat]:
    """Pull number+label+source out of a list of fact-shaped strings."""
    stats: list[KpiStat] = []
    for line in lines:
        nums = _extract_numbers(line)
        if not nums:
            continue
        value = nums[0][:15]
        label = _strip_citations(line)
        # Trim label so the stat card stays readable
        if value in label:
            label = label.replace(value, "", 1).strip(" :,-—–")
        label = label[:60] or "Metric"
        note = _first_citation(line)
        stats.append(KpiStat(value=value, label=label, note=f"[{note}]" if note else None))
        if len(stats) >= max_n:
            break
    return stats


def _deterministic_design(outline: dict) -> DeckSpec:
    """No-LLM fallback if the Designer call fails or returns junk.

    Maps each outline item to a sensible layout using the documented rules
    AND scans the actual content (not just role labels) to promote items
    into kpi_stats / comparison / quote when the data justifies it.
    """
    title = outline.get("title", "Untitled Deck")
    subtitle = outline.get("subtitle")
    items = outline.get("outline") or []
    slides: list[SlideSpec] = []

    for i, item in enumerate(items):
        role = (item.get("narrative_role") or "").lower()
        headline = item.get("headline") or item.get("title") or ""
        key_points = list(item.get("key_points") or [])
        supporting = list(item.get("supporting_data") or [])

        # Fixed-role slides first.
        if role == "cover" or i == 0:
            slides.append(SlideSpec(layout="cover", title=title, subtitle=subtitle))
            continue
        if role == "conclusion" or i == len(items) - 1:
            slides.append(SlideSpec(layout="closing", title="Q&A",
                                    subtitle=headline or None))
            continue

        # Quote takes precedence if drafter flagged one.
        if role == "quote" or item.get("quotable"):
            q = item.get("quotable") or _looks_like_quote(" ".join(supporting))
            if q:
                slides.append(SlideSpec(
                    layout="quote",
                    title=headline,
                    quote=q,
                    attribution=item.get("quotable_source") or _first_citation(" ".join(supporting)) or "",
                ))
                continue

        # Comparison if drafter set both columns, OR if headline has 'vs' / 'versus'.
        if (role == "comparison" and item.get("compare_a") and item.get("compare_b")) or \
           re.search(r"\b(vs|versus)\b", headline, re.IGNORECASE):
            ca = item.get("compare_a") or _split_vs(headline)[0]
            cb = item.get("compare_b") or _split_vs(headline)[1]
            half = max(1, len(key_points) // 2)
            slides.append(SlideSpec(
                layout="comparison",
                title=headline,
                columns=[
                    ComparisonColumn(heading=str(ca), bullets=key_points[:half]),
                    ComparisonColumn(heading=str(cb), bullets=key_points[half:] or supporting[:3]),
                ],
            ))
            continue

        # kpi_stats — auto-detect: if supporting_data or key_points have 2+ number-shaped lines,
        # render as stat cards regardless of role label.
        candidates = supporting + key_points
        numeric_lines = [l for l in candidates if _extract_numbers(l)]
        if (role == "metrics" or len(numeric_lines) >= 2) and len(numeric_lines) >= 2:
            stats = _kpi_from_lines(numeric_lines, max_n=4)
            if stats:
                slides.append(SlideSpec(layout="kpi_stats", title=headline, stats=stats))
                continue

        # Default — bullets. Combine key_points + any unused supporting_data as bullets.
        bullets = list(key_points) or list(supporting)
        slides.append(SlideSpec(layout="bullets", title=headline, bullets=bullets))

    # Insert section dividers between narrative phase changes.
    # Heuristic: section divider before slides 3, 6, 9 (so a 10-slide deck gets
    # 3 dividers, not just one). Skip if neighbour is already a section/cover/closing.
    final: list[SlideSpec] = []
    section_at = {3, 6, 9}
    for i, s in enumerate(slides):
        if i in section_at and s.layout not in ("cover", "section", "closing"):
            divider_title = s.title or "Section"
            final.append(SlideSpec(layout="section", title=divider_title))
        final.append(s)
    return DeckSpec(title=title, subtitle=subtitle, slides=final)


def _split_vs(text: str) -> tuple[str, str]:
    parts = re.split(r"\s+(?:vs|versus)\s+", text or "", maxsplit=1, flags=re.IGNORECASE)
    return (parts[0].strip(), parts[1].strip() if len(parts) > 1 else "Alternative")


def _extract_json_object(text: str) -> dict | None:
    """Robust JSON-object extraction.

    Qwen occasionally emits stray characters (markdown fences, leading `[`,
    trailing commentary). Strategy: strip code fences, fix the known
    `{[` artefact (where Qwen interleaves an array literal right after the
    opening brace), find the first `{`, then locate its matching `}` by
    tracking brace depth (respecting strings + escapes). Fall back to a
    plain json.loads as last resort.
    """
    s = re.sub(r"^```(?:json)?\s*|\s*```$", "", text.strip(), flags=re.DOTALL)
    # Repair common Qwen artefacts.
    # Case A: "{[]" (empty bracket between { and key)        →   {
    # Case B: "{[{..."  (array-of-objects wrapped in object) →   try as array
    # Case C: "{[\"key" or "{[\nkey" (stray [ before key)    →   strip the [
    s = re.sub(r"\{\s*\[\s*\]\s*", "{", s)
    # Case B handled BEFORE case C: if the bracket opens an array of objects,
    # peel both wrappers, parse as a list. Then disambiguate:
    #   - elements look like deck specs (have "title"/"slides")  → return arr[0]
    #   - elements look like slides (have "layout")              → wrap as {"slides": arr}
    wrapped = re.match(r"^\s*\{\s*\[\s*(?=\{)", s)
    if wrapped:
        inner = re.sub(r"^\s*\{\s*\[\s*", "[", s)
        inner = re.sub(r"\s*\]\s*\}\s*$", "]", inner)
        try:
            arr = json.loads(inner)
            if isinstance(arr, list) and arr and isinstance(arr[0], dict):
                first_keys = set(arr[0].keys())
                if "layout" in first_keys:
                    return {"slides": arr}
                return arr[0]
        except Exception:
            # fall through to standard brace-matching extraction
            pass
    s = re.sub(r"\{\s*\[(\s*[\"\n])", r"{\1", s)         # {["key" or {[\n → {key or {\n
    start = s.find("{")
    if start == -1:
        return None
    depth = 0
    in_str = False
    esc = False
    for i in range(start, len(s)):
        ch = s[i]
        if in_str:
            if esc:
                esc = False
            elif ch == "\\":
                esc = True
            elif ch == '"':
                in_str = False
            continue
        if ch == '"':
            in_str = True
        elif ch == "{":
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0:
                try:
                    return json.loads(s[start : i + 1])
                except Exception:
                    return None
    try:
        return json.loads(s)
    except Exception:
        return None


# ----- python-pptx renderer -----
# Brooker Group brand palette (sourced from the Dec-2021 template's theme XML).
# Red is the primary brand colour; navy is the secondary; grey is muted text.
_PRIMARY  = RGBColor(0xEE, 0x31, 0x35)   # Brooker red       (accent1)
_DEEP_RED = RGBColor(0xAE, 0x13, 0x2A)   # deep red          (accent2)
_ACCENT   = RGBColor(0x00, 0x28, 0x56)   # navy              (accent4)
_INK      = RGBColor(0x00, 0x00, 0x00)   # black text        (dk1)
_MUTED    = RGBColor(0x5F, 0x5F, 0x5F)   # secondary text    (dk2)
_BAND     = RGBColor(0xF2, 0xF2, 0xF2)   # very light grey band background


def _slide_dims(prs):
    return prs.slide_width, prs.slide_height


def _add_textbox(slide, left, top, width, height, *, text="", size=18,
                 bold=False, color=None, align=None, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        if align is not None:
            p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            if color is not None:
                r.font.color.rgb = color
    return box


def _add_band(slide, left, top, width, height, fill_color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    return shp


def _pick_layout(prs, *, want_body: bool):
    """Return a slide_layout that has a body placeholder when want_body is True.

    The IC template ships with many custom layouts whose body placeholders
    live at non-standard `placeholder_format.idx` values. We probe each layout
    and pick the first one whose layout-level placeholders include at least
    one non-title placeholder (idx != 0).
    """
    if not want_body:
        return prs.slide_layouts[0]
    for layout in prs.slide_layouts:
        idxs = [ph.placeholder_format.idx for ph in layout.placeholders]
        # idx 0 is the title; any other idx means a content/body placeholder.
        if any(i != 0 for i in idxs):
            return layout
    # No layout with a body — fall back to layout 1 if present, else 0.
    return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]


def _find_body_placeholder(slide):
    """Find a placeholder we can write bullets into, excluding the title."""
    title_ph = slide.shapes.title
    for ph in slide.placeholders:
        if ph == title_ph:
            continue
        # any non-title placeholder with a text_frame is usable
        if ph.has_text_frame:
            return ph
    return None


def _fit_title(slide, headline: str) -> tuple[float, float]:
    """Set the slide title, shrinking the font when the headline is long so it
    doesn't wrap into the body area. Returns (title_bottom_y_inches, used_pt)
    so the body textbox can be placed below it if needed."""
    title_ph = slide.shapes.title
    if title_ph is None or not headline:
        return (0.0, 32.0)
    title_ph.text = headline

    # Length-bucketed font sizing — keeps the title to ~2 lines max.
    length = len(headline)
    if length <= 35:
        size_pt = 32.0
    elif length <= 55:
        size_pt = 26.0
    elif length <= 80:
        size_pt = 22.0
    else:
        size_pt = 18.0
    # Apply to all runs (template might default to 40pt+).
    tf = title_ph.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size_pt)
    # Best-effort bottom-y estimate (EMU → inches).
    bottom_in = ((title_ph.top or 0) + (title_ph.height or 0)) / 914400
    return (bottom_in, size_pt)


def _write_bullets(slide, bullets: list[str]) -> None:
    """Write bullets. If a body placeholder exists it goes there; otherwise we
    add a textbox positioned BELOW the (possibly-tall) title so long titles
    don't overlap the bullets."""
    if not bullets:
        return
    body = _find_body_placeholder(slide)
    if body is None:
        # No body placeholder in this layout — draw our own textbox.
        # Anchor below the title's bottom so a 2-line title doesn't clobber us.
        title_ph = slide.shapes.title
        try:
            title_bottom = ((title_ph.top or 0) + (title_ph.height or 0)) / 914400
        except Exception:
            title_bottom = 1.5
        prs = slide.part.package.presentation_part.presentation
        slide_w = prs.slide_width / 914400
        slide_h = prs.slide_height / 914400
        top_in = max(title_bottom + 0.15, 1.4)
        body = slide.shapes.add_textbox(
            Inches(0.6), Inches(top_in),
            Inches(slide_w - 1.2), Inches(slide_h - top_in - 0.5),
        )
        tf = body.text_frame
        tf.word_wrap = True
    else:
        tf = body.text_frame
    tf.text = bullets[0]
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
    # Bullet font sizing — shrink when there are many bullets so long key_points
    # don't overflow the placeholder.
    n_bullets = len(bullets)
    avg_len = sum(len(b) for b in bullets) / max(1, n_bullets)
    if n_bullets >= 6 or avg_len > 80:
        bullet_size = Pt(14)
    elif n_bullets >= 4 or avg_len > 50:
        bullet_size = Pt(16)
    else:
        bullet_size = Pt(18)
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = bullet_size


def _find_layout_by_name(prs, *names: str):
    """Return the first slide_layout whose .name matches any of `names`
    (case-insensitive, substring match)."""
    lower = [n.lower() for n in names]
    for layout in prs.slide_layouts:
        if any(want in (layout.name or "").lower() for want in lower):
            return layout
    return None


def _render_cover(prs, fallback_layout, spec_title: str, subtitle: str | None) -> None:
    """Use the template's native Title Slide layout if present; else
    fall back to a hand-drawn cover with brand accents."""
    native = _find_layout_by_name(prs, "title slide", "cover")
    if native is not None:
        slide = prs.slides.add_slide(native)
        if slide.shapes.title is not None:
            slide.shapes.title.text = spec_title
        # Fill any subtitle-shaped placeholder (idx==1, or shape name contains
        # "subtitle"). Skip the title placeholder.
        title_ph = slide.shapes.title
        for ph in slide.placeholders:
            if ph == title_ph:
                continue
            name = (ph.name or "").lower()
            if subtitle and ("subtitle" in name or ph.placeholder_format.idx == 1):
                ph.text_frame.text = subtitle
                break
        return

    # Fallback — template has no Title Slide layout.
    slide = prs.slides.add_slide(fallback_layout)
    w, h = _slide_dims(prs)
    _add_band(slide, 0, 0, w, Inches(0.6), _PRIMARY)
    _add_textbox(slide, Inches(0.6), Inches(2.3), w - Inches(1.2), Inches(2.0),
                 text=spec_title, size=44, bold=True, color=_INK)
    if subtitle:
        _add_textbox(slide, Inches(0.6), Inches(4.4), w - Inches(1.2), Inches(0.8),
                     text=subtitle, size=22, color=_ACCENT, italic=True)


def _render_section(prs, blank_layout, title: str) -> None:
    slide = prs.slides.add_slide(blank_layout)
    w, h = _slide_dims(prs)
    _add_band(slide, 0, Inches(2.5), w, Inches(2.5), _PRIMARY)
    _add_textbox(slide, Inches(0.6), Inches(3.0), w - Inches(1.2), Inches(1.5),
                 text=title, size=40, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))


def _render_bullets(prs, body_layout, s: SlideSpec) -> None:
    slide = prs.slides.add_slide(body_layout)
    # Use the title-fitter so long headlines don't blow out the placeholder.
    _fit_title(slide, s.title or "")
    _write_bullets(slide, list(s.bullets))


def _render_comparison(prs, blank_layout, s: SlideSpec) -> None:
    slide = prs.slides.add_slide(blank_layout)
    w, h = _slide_dims(prs)
    # Title
    _add_textbox(slide, Inches(0.6), Inches(0.4), w - Inches(1.2), Inches(0.8),
                 text=s.title or "", size=28, bold=True, color=_INK)
    cols = (s.columns or [])[:2]
    if not cols:
        return
    gutter = Inches(0.4)
    col_w = (w - Inches(1.2) - gutter) / 2
    top = Inches(1.5)
    col_h = h - Inches(1.9)
    for i, col in enumerate(cols):
        left = Inches(0.6) + (col_w + gutter) * i
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, col_h)
        card.fill.solid()
        card.fill.fore_color.rgb = _BAND
        card.line.color.rgb = _PRIMARY
        card.line.width = Pt(1)
        # Heading
        _add_textbox(slide, left + Inches(0.3), top + Inches(0.2),
                     col_w - Inches(0.6), Inches(0.7),
                     text=col.heading or "", size=20, bold=True, color=_PRIMARY)
        # Bullets
        bbox = slide.shapes.add_textbox(
            left + Inches(0.3), top + Inches(1.0),
            col_w - Inches(0.6), col_h - Inches(1.2),
        )
        tf = bbox.text_frame
        tf.word_wrap = True
        bs = col.bullets or []
        if not bs:
            tf.text = ""
        else:
            tf.text = "• " + bs[0]
            for b in bs[1:]:
                p = tf.add_paragraph()
                p.text = "• " + b
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(16)
                r.font.color.rgb = _INK


def _render_kpi_stats(prs, blank_layout, s: SlideSpec) -> None:
    slide = prs.slides.add_slide(blank_layout)
    w, h = _slide_dims(prs)
    _add_textbox(slide, Inches(0.6), Inches(0.4), w - Inches(1.2), Inches(0.8),
                 text=s.title or "", size=28, bold=True, color=_INK)
    stats = (s.stats or [])[:4]
    if not stats:
        return
    gutter = Inches(0.3)
    card_w = (w - Inches(1.2) - gutter * (len(stats) - 1)) / len(stats)
    top = Inches(2.2)
    card_h = Inches(3.6)
    for i, stat in enumerate(stats):
        left = Inches(0.6) + (card_w + gutter) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = _BAND
        card.line.color.rgb = _ACCENT
        card.line.width = Pt(1)
        # Big value
        _add_textbox(slide, left, top + Inches(0.4), card_w, Inches(1.6),
                     text=stat.value or "", size=44, bold=True,
                     color=_PRIMARY, align=PP_ALIGN.CENTER)
        # Label
        _add_textbox(slide, left + Inches(0.2), top + Inches(2.0),
                     card_w - Inches(0.4), Inches(0.8),
                     text=stat.label or "", size=14, bold=True,
                     color=_INK, align=PP_ALIGN.CENTER)
        # Note
        if stat.note:
            _add_textbox(slide, left + Inches(0.2), top + Inches(2.8),
                         card_w - Inches(0.4), Inches(0.6),
                         text=stat.note, size=11, color=_MUTED,
                         italic=True, align=PP_ALIGN.CENTER)


def _render_quote(prs, blank_layout, s: SlideSpec) -> None:
    slide = prs.slides.add_slide(blank_layout)
    w, h = _slide_dims(prs)
    # Left accent bar
    _add_band(slide, Inches(0.8), Inches(2.0), Inches(0.15), Inches(3.5), _PRIMARY)
    quote_text = (s.quote or "").strip()
    if not quote_text.startswith("“"):
        quote_text = f"“{quote_text}”"
    _add_textbox(slide, Inches(1.3), Inches(2.0), w - Inches(2.0), Inches(3.0),
                 text=quote_text, size=28, italic=True, color=_INK)
    if s.attribution:
        _add_textbox(slide, Inches(1.3), Inches(5.0), w - Inches(2.0), Inches(0.6),
                     text=f"— {s.attribution}", size=16, color=_MUTED)


def _render_closing(prs, blank_layout, s: SlideSpec) -> None:
    slide = prs.slides.add_slide(blank_layout)
    w, h = _slide_dims(prs)
    title = s.title or "Q&A"
    _add_textbox(slide, Inches(0.6), Inches(2.8), w - Inches(1.2), Inches(1.5),
                 text=title, size=54, bold=True, color=_PRIMARY, align=PP_ALIGN.CENTER)
    if s.subtitle:
        _add_textbox(slide, Inches(0.6), Inches(4.4), w - Inches(1.2), Inches(0.8),
                     text=s.subtitle, size=20, color=_MUTED,
                     italic=True, align=PP_ALIGN.CENTER)


def _render(spec: DeckSpec, dept_id: str = "ic") -> Path:
    template = Path(_template_for(dept_id))
    log.info("render.template", dept_id=dept_id, template=template.name, exists=template.is_file())
    if template.is_file():
        prs = Presentation(str(template))
        # Clear any pre-existing slides from the template.
        sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
        for sldId in list(sldIdLst):
            rId = sldId.rId
            prs.part.drop_rel(rId)
            sldIdLst.remove(sldId)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # Layouts: blank (idx 6 in Office defaults) when we want full control;
    # bullet layout when we want Office's built-in title+body styling.
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else _pick_layout(prs, want_body=False)
    body_layout = _pick_layout(prs, want_body=True)

    # Always start with a cover even if the spec didn't ask for one.
    first = spec.slides[0] if spec.slides else None
    if not first or first.layout != "cover":
        _render_cover(prs, blank_layout, spec.title, spec.subtitle)
        content_slides = list(spec.slides)
    else:
        _render_cover(prs, blank_layout,
                      first.title or spec.title,
                      first.subtitle or spec.subtitle)
        content_slides = list(spec.slides[1:])

    dispatch = {
        "cover":     lambda s: _render_cover(prs, blank_layout, s.title or spec.title, s.subtitle),
        "section":   lambda s: _render_section(prs, blank_layout, s.title or ""),
        "bullets":   lambda s: _render_bullets(prs, body_layout, s),
        "comparison":lambda s: _render_comparison(prs, blank_layout, s),
        "kpi_stats": lambda s: _render_kpi_stats(prs, blank_layout, s),
        "quote":     lambda s: _render_quote(prs, blank_layout, s),
        "closing":   lambda s: _render_closing(prs, blank_layout, s),
    }
    # Fallback so we never produce a 1-slide deck. If the LLM returned no
    # content slides, render a single bullets slide noting the failure mode.
    if not content_slides:
        content_slides = [SlideSpec(
            layout="bullets",
            title="Context limited",
            bullets=[
                "The model produced a cover but no content slides.",
                "This usually means the LLM JSON was incomplete or the retrieved context was too thin.",
                "Try rephrasing the brief, or attach a source document.",
            ],
        )]
    for s in content_slides:
        renderer = dispatch.get(s.layout, dispatch["bullets"])
        try:
            renderer(s)
        except Exception as exc:
            log.warning("render_layout_failed", layout=s.layout, err=str(exc))
            _render_bullets(prs, body_layout, s)   # safest fallback
        if s.notes:
            try:
                prs.slides[-1].notes_slide.notes_text_frame.text = s.notes
            except Exception:
                pass

    fname = f"deck_{datetime.now(timezone.utc).strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:6]}.pptx"
    out = OUTPUT_DIR / fname
    prs.save(str(out))
    return out


# ----- rich-content helpers -----

def _build_image_source(embed: ImageEmbed) -> ImageSource:
    """Convert an ImageEmbed Pydantic model into an ImageSource dict."""
    if embed.url is not None:
        return {"url": embed.url}
    if embed.path is not None:
        return {"path": embed.path}
    if embed.minio_key is not None:
        return {"minio_key": embed.minio_key}
    raise ValueError("ImageEmbed must have exactly one of: url, path, minio_key")


def _embed_rich_content_docx(
    doc,
    images: list[ImageEmbed],
    charts: list[ChartEmbed],
    mermaid_items: list[MermaidEmbed],
) -> None:
    """Embed images, charts, and Mermaid diagrams into an open python-docx Document.

    All insertions happen AFTER the document body is already built.  Failures on
    individual items are logged and skipped — we never let a bad image kill the
    whole report.
    """
    for img in images:
        try:
            src = _build_image_source(img)
            result = embed_image_in_docx(
                doc, src,
                section_hint=img.section_hint,
                caption=img.caption,
                width_inches=img.width_inches,
            )
            log.info("report.image_embedded", section_hint=img.section_hint, result=result)
        except Exception as exc:
            log.warning("report.image_embed_failed", error=str(exc))

    for chart in charts:
        try:
            validated_spec = validate_chart_spec(chart.spec)
            png_bytes = render_chart_png(validated_spec)
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp.write(png_bytes)
                tmp_path = tmp.name
            try:
                result = embed_image_in_docx(
                    doc, {"path": tmp_path},
                    section_hint=chart.section_hint,
                    caption=chart.caption,
                    width_inches=chart.width_inches,
                )
                log.info("report.chart_embedded", section_hint=chart.section_hint, result=result)
            finally:
                Path(tmp_path).unlink(missing_ok=True)
        except Exception as exc:
            log.warning("report.chart_embed_failed", error=str(exc))

    for item in mermaid_items:
        try:
            png_bytes = render_mermaid_png(
                item.mermaid,
                theme=item.theme,  # type: ignore[arg-type]
                background=item.background,
            )
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp.write(png_bytes)
                tmp_path = tmp.name
            try:
                result = embed_image_in_docx(
                    doc, {"path": tmp_path},
                    section_hint=item.section_hint,
                    caption=item.caption,
                    width_inches=item.width_inches,
                )
                log.info("report.mermaid_embedded", section_hint=item.section_hint,
                         result=result)
            finally:
                Path(tmp_path).unlink(missing_ok=True)
        except Exception as exc:
            log.warning("report.mermaid_embed_failed", error=str(exc))


def _embed_rich_content_pptx(
    prs,
    images: list[ImageEmbed],
    charts: list[ChartEmbed],
    mermaid_items: list[MermaidEmbed],
) -> None:
    """Embed images, charts, and Mermaid diagrams into an open python-pptx Presentation.

    All insertions happen AFTER the deck is rendered.  Failures on individual
    items are logged and skipped — we never let a bad chart kill the whole deck.
    """
    for img in images:
        try:
            src = _build_image_source(img)
            result = embed_image_in_pptx(
                prs, src,
                slide_index=img.slide_index,
                slide_title_hint=img.slide_title_hint,
                caption=img.caption,
                width_inches=img.width_inches,
            )
            log.info("compose.image_embedded",
                     slide_index=img.slide_index,
                     slide_title_hint=img.slide_title_hint,
                     result=result)
        except Exception as exc:
            log.warning("compose.image_embed_failed", error=str(exc))

    for chart in charts:
        try:
            validated_spec = validate_chart_spec(chart.spec)
            png_bytes = render_chart_png(validated_spec)
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp.write(png_bytes)
                tmp_path = tmp.name
            try:
                result = embed_image_in_pptx(
                    prs, {"path": tmp_path},
                    slide_index=chart.slide_index,
                    slide_title_hint=chart.slide_title_hint,
                    caption=chart.caption,
                    width_inches=chart.width_inches,
                )
                log.info("compose.chart_embedded",
                         slide_index=chart.slide_index,
                         slide_title_hint=chart.slide_title_hint,
                         result=result)
            finally:
                Path(tmp_path).unlink(missing_ok=True)
        except Exception as exc:
            log.warning("compose.chart_embed_failed", error=str(exc))

    for item in mermaid_items:
        try:
            png_bytes = render_mermaid_png(
                item.mermaid,
                theme=item.theme,  # type: ignore[arg-type]
                background=item.background,
            )
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp.write(png_bytes)
                tmp_path = tmp.name
            try:
                result = embed_image_in_pptx(
                    prs, {"path": tmp_path},
                    slide_index=item.slide_index,
                    slide_title_hint=item.slide_title_hint,
                    caption=item.caption,
                    width_inches=item.width_inches,
                )
                log.info("compose.mermaid_embedded",
                         slide_index=item.slide_index,
                         slide_title_hint=item.slide_title_hint,
                         result=result)
            finally:
                Path(tmp_path).unlink(missing_ok=True)
        except Exception as exc:
            log.warning("compose.mermaid_embed_failed", error=str(exc))


def _extract_tables_from_outline(outline: dict) -> tuple[dict, list[dict]]:
    """Scan all text fields in the drafter outline for ```table blocks.

    Modifies a shallow copy of the outline (does not mutate the original).
    Returns (cleaned_outline, [table_specs]).
    """
    all_table_specs: list[dict] = []

    def _clean(text: str) -> str:
        nonlocal all_table_specs
        cleaned, specs = extract_tables_from_text(text)
        all_table_specs.extend(specs)
        return cleaned

    cleaned_outline = dict(outline)
    cleaned_items: list[dict] = []
    for item in (outline.get("outline") or []):
        cleaned_item = dict(item)
        if "headline" in cleaned_item:
            cleaned_item["headline"] = _clean(cleaned_item["headline"])
        if "key_points" in cleaned_item:
            cleaned_item["key_points"] = [_clean(kp) for kp in (cleaned_item["key_points"] or [])]
        if "supporting_data" in cleaned_item:
            cleaned_item["supporting_data"] = [
                _clean(sd) for sd in (cleaned_item["supporting_data"] or [])
            ]
        cleaned_items.append(cleaned_item)

    cleaned_outline["outline"] = cleaned_items
    return cleaned_outline, all_table_specs


# ----- endpoint -----
@app.post("/compose", response_model=ComposeResponse)
async def compose(req: ComposeRequest) -> ComposeResponse:
    t0 = time.monotonic()
    # Accept both `brief` and the legacy `query` field from slack-bot router.
    brief = req.brief or req.query or ""
    if not brief.strip():
        raise HTTPException(400, "brief (or query) is required")

    log.info("compose.start", brief=brief[:120], dept_id=req.dept_id)
    hits = await _retrieve(brief, req.dept_id)
    context = "\n\n".join(f"[{h['source']}] {h['excerpt']}" for h in hits[:8])
    # Two-agent pipeline: Drafter → Designer → Renderer.
    raw_outline = await _draft_outline_with_retry(brief, context, req.title, req.audience)

    # Strip ```table blocks from outline text before the designer sees it;
    # extracted specs are inserted as table slides after the deck is rendered.
    outline, deck_table_specs = _extract_tables_from_outline(raw_outline)

    spec = await _design_layouts(outline)
    out_path = _render(spec, req.dept_id)

    # ----- rich-content insertion (runs AFTER the pptx is rendered) -----
    need_pptx_save = bool(req.images or req.charts or req.mermaid or deck_table_specs)
    if need_pptx_save:
        from pptx import Presentation as _Prs
        prs = _Prs(str(out_path))
        # Insert tables extracted from the drafter outline as slides.
        for tspec in deck_table_specs:
            try:
                validated_tspec = validate_table_spec(tspec)
                add_table_to_pptx(prs, validated_tspec)
                log.info("compose.table_slide_inserted", title=tspec.get("title", ""))
            except Exception as exc:
                log.warning("compose.table_slide_failed", error=str(exc))
        # Insert explicitly requested rich-content embeds.
        if req.images or req.charts or req.mermaid:
            _embed_rich_content_pptx(prs, req.images, req.charts, req.mermaid)
            log.info("compose.rich_content_done",
                     images=len(req.images), charts=len(req.charts),
                     mermaid=len(req.mermaid))
        prs.save(str(out_path))

    elapsed_ms = int((time.monotonic() - t0) * 1000)
    log.info(
        "compose.done",
        file=out_path.name,
        slides=len(spec.slides),
        sources=len(hits),
        elapsed_ms=elapsed_ms,
    )

    file_url = f"http://deck-writer:3050/files/{out_path.name}"
    answer = (
        f"Drafted *{spec.title}* — {len(spec.slides)} slides "
        f"({len(hits)} sources). File: `{out_path.name}`"
    )
    return ComposeResponse(
        answer=answer,
        confidence="Medium",
        file_path=str(out_path),
        file_name=out_path.name,
        file_url=file_url,
        sources=[
            {"source": h["source"], "excerpt": h["excerpt"][:120], "score": h["score"]}
            for h in hits[:8]
        ],
    )


# Slack-bot router compatibility: when it POSTs to /query, just forward to /compose.
@app.post("/query", response_model=ComposeResponse)
async def query_alias(req: ComposeRequest) -> ComposeResponse:
    return await compose(req)


# --------------------------------------------------------------------------
# /report — same retrieval + drafter step as /compose, but emits a .docx
# review document instead of a pptx. The Drafter outline becomes the
# document's section structure; each section gets a heading + paragraphs +
# bulleted key points + a Sources block at the end.
# --------------------------------------------------------------------------
_REPORT_OUTPUT_DIR = OUTPUT_DIR.parent / "reports"
_REPORT_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
_XLSX_OUTPUT_DIR = _REPORT_OUTPUT_DIR  # reuse /data/reports for xlsx output


def _render_docx(outline: dict, hits: list[dict], *,
                  extra_table_specs: list[dict] | None = None) -> Path:
    """Render a .docx from a drafter outline.

    extra_table_specs: pre-parsed TableSpec dicts extracted from drafter text
    via extract_tables_from_text(); inserted at the end of the document body
    before the Sources section.
    """
    from docx import Document

    doc = Document()
    title = outline.get("title", "Untitled Report")
    subtitle = outline.get("subtitle")

    doc.add_heading(title, level=0)
    if subtitle:
        sub = doc.add_paragraph(subtitle)
        sub.runs[0].italic = True

    # Executive summary built from headlines of all non-cover/non-conclusion items
    items = outline.get("outline") or []
    summary_pts = [
        i.get("headline") for i in items
        if (i.get("narrative_role") or "").lower() not in ("cover", "conclusion")
        and i.get("headline")
    ]
    if summary_pts:
        doc.add_heading("Executive summary", level=1)
        for pt in summary_pts[:6]:
            doc.add_paragraph(pt, style="List Bullet")

    # Per-section content
    for item in items:
        role = (item.get("narrative_role") or "").lower()
        if role in ("cover", "conclusion"):
            continue
        head = item.get("headline") or "Section"
        doc.add_heading(head, level=1)
        for kp in (item.get("key_points") or []):
            doc.add_paragraph(kp, style="List Bullet")
        supporting = item.get("supporting_data") or []
        if supporting:
            doc.add_heading("Supporting data", level=2)
            for s in supporting:
                doc.add_paragraph(s, style="List Bullet")

    # Tables extracted from drafter text — appended after section content
    for tspec in (extra_table_specs or []):
        try:
            validated = validate_table_spec(tspec)
            add_table_to_docx(doc, validated)
            log.info("report.table_inserted", title=tspec.get("title", ""))
        except Exception as exc:
            log.warning("report.table_insert_failed", error=str(exc))

    # Sources section
    if hits:
        doc.add_heading("Sources", level=1)
        seen: set[str] = set()
        for h in hits[:20]:
            src = h.get("source")
            if not src or src in seen:
                continue
            seen.add(src)
            doc.add_paragraph(f"• {src}")

    fname = f"report_{datetime.now(timezone.utc).strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:6]}.docx"
    out = _REPORT_OUTPUT_DIR / fname
    doc.save(str(out))
    return out


@app.post("/report")
async def report(req: ComposeRequest) -> ComposeResponse:
    """Generate a .docx review document from the same drafter outline used for decks.

    Optional rich-content fields on ComposeRequest:
      images  — list[ImageEmbed]: PNG/JPEG images inserted after the docx is built
      charts  — list[ChartEmbed]: matplotlib charts rendered to PNG then embedded
      mermaid — list[MermaidEmbed]: Mermaid diagrams rendered to PNG then embedded

    Table extraction: when the drafter LLM emits ```table JSON blocks, they are
    stripped from the narrative text and rendered as styled Word tables in the
    document body via add_table_to_docx.
    """
    brief = (req.brief or req.query or "").strip()
    if not brief:
        raise HTTPException(400, "brief (or query) is required")
    log.info("report.start", brief=brief[:120], dept_id=req.dept_id)
    hits = await _retrieve(brief, req.dept_id)
    context = "\n\n".join(f"[{h['source']}] {h['excerpt']}" for h in hits[:8])
    raw_outline = await _draft_outline_with_retry(brief, context, req.title, req.audience)

    # Extract ```table blocks from all free-text fields in the outline so they
    # render as structured Word tables rather than raw JSON in the narrative.
    outline, table_specs = _extract_tables_from_outline(raw_outline)

    out_path = _render_docx(outline, hits, extra_table_specs=table_specs)

    # ----- rich-content insertion (runs AFTER the docx is rendered) -----
    if req.images or req.charts or req.mermaid:
        from docx import Document as _Doc
        doc = _Doc(str(out_path))
        _embed_rich_content_docx(doc, req.images, req.charts, req.mermaid)
        doc.save(str(out_path))
        log.info("report.rich_content_done",
                 images=len(req.images), charts=len(req.charts),
                 mermaid=len(req.mermaid))

    fname = out_path.name
    log.info("report.done", file=fname, sections=len(outline.get("outline") or []),
             sources=len(hits), tables=len(table_specs))
    return ComposeResponse(
        answer=f"Drafted *{outline.get('title','Report')}* — "
               f"{len(outline.get('outline') or [])} sections, {len(hits)} sources. "
               f"File: `{fname}`",
        confidence="Medium",
        file_path=str(out_path),
        file_name=fname,
        file_url=f"http://deck-writer:3050/reports/{fname}",
        sources=[
            {"source": h["source"], "excerpt": h["excerpt"][:120], "score": h["score"]}
            for h in hits[:8]
        ],
    )


@app.get("/reports/{name}")
async def get_report(name: str) -> FileResponse:
    safe = re.sub(r"[^A-Za-z0-9_.\-]", "", name)
    p = _REPORT_OUTPUT_DIR / safe
    if not p.is_file():
        raise HTTPException(404, "report not found")
    return FileResponse(
        path=str(p),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        filename=safe,
    )


# --------------------------------------------------------------------------
# /compose-xlsx — free-form Excel workbook builder
# --------------------------------------------------------------------------
@app.post("/compose-xlsx")
async def compose_xlsx_endpoint(req: ComposeXlsxRequest) -> FileResponse:
    """Build a .xlsx workbook from an XlsxComposeSpec.

    Request body:
        spec        — XlsxComposeSpec dict (sheets, optional charts, optional metadata)
        filename    — output filename, e.g. "capital_report.xlsx" (basename only)
        caller_dept — requesting department; used to scope the output directory

    Returns the .xlsx file as an attachment.
    """
    # Sanitise filename — strip path separators so callers can't traverse dirs.
    safe_filename = re.sub(r"[^A-Za-z0-9_.\-]", "_", req.filename) or "report.xlsx"
    if not safe_filename.endswith(".xlsx"):
        safe_filename += ".xlsx"

    # Validate spec eagerly so we return a 400 before writing any file.
    try:
        validated_spec = validate_xlsx_spec(req.spec)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    dept = re.sub(r"[^A-Za-z0-9_\-]", "", req.caller_dept or "shared")
    out_dir = Path("/data/reports") / dept
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / safe_filename

    try:
        result = compose_xlsx(validated_spec, out_path)
    except (ValueError, OSError) as exc:
        log.error("compose_xlsx.failed", error=str(exc), caller_dept=dept)
        raise HTTPException(status_code=500, detail=str(exc)) from exc

    log.info("compose_xlsx.done",
             out=result["out"],
             sheets=result["sheets"],
             cells_written=result["cells_written"],
             charts=result["charts"],
             caller_dept=dept)

    return FileResponse(
        path=str(out_path),
        filename=safe_filename,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


# --------------------------------------------------------------------------
# /report/cac-meeting — DETERMINISTIC CAC committee report from the live Excel
# Online data pack (MS Graph). Numbers are read straight from the workbook cells
# (no LLM, so figures can never be altered/rounded by a model); limit breaches
# are computed in code. Returns the .docx via file_url so the slack-bot uploads
# it to the thread instead of posting markdown. Mirrors the deterministic builder
# in scripts/gen_cac_report_from_xlsx.py via services/shared/cac_report_docx.py.
# --------------------------------------------------------------------------
@app.get("/report/cac-meeting")
async def cac_meeting_report(
    share_url: str | None = None,
    caller_dept: str | None = None,
) -> dict:
    """Build the CAC committee report .docx from the live Data Pack.

    Department guard (defence in depth — slack-bot also gates routing):
      The CAC report is a CAC-owned artefact. Only the CAC department may invoke
      this endpoint over the Slack pipeline. `caller_dept` is the channel-dept
      slack-bot resolved for the message; it is omitted for direct curl/test calls
      (which are allowed so the endpoint stays diagnosable).
    """
    from datetime import date

    from services.shared.cac_report_docx import (
        build_cac_report_docx,
        pack_from_workbook,
    )
    from services.shared.ms_graph_excel import GraphExcel

    # caller_dept None == direct call (curl/test); anything else must be 'cac'.
    if caller_dept is not None and caller_dept != "cac":
        log.warning("cac_meeting_report.dept_blocked", caller_dept=caller_dept)
        raise HTTPException(
            403,
            f"The CAC meeting report can only be requested from a CAC context "
            f"(caller_dept='{caller_dept}'). Post the request in #cac-committee.",
        )

    share_url = (share_url or os.getenv("CAC_DATA_PACK_SHARE_URL", "")).strip()
    if not share_url:
        raise HTTPException(
            422,
            "No Excel Online link given. Paste a SharePoint/OneDrive share link "
            "after the command (e.g. `[cac-report] https://...`) or set "
            "CAC_DATA_PACK_SHARE_URL.",
        )
    if not share_url.lower().startswith(("http://", "https://")):
        raise HTTPException(422, f"Not a valid share link: {share_url[:80]!r}")

    try:
        gx = GraphExcel.from_env()
        workbook = await gx.read_workbook_by_share_url(share_url)
    except Exception as exc:  # noqa: BLE001
        log.error("cac_meeting_report.excel_read_failed", error=str(exc))
        raise HTTPException(502, f"Could not read the Excel Online data pack: {exc}")

    month = date.today().strftime("%B %Y")
    fname = f"CAC_Report_{month.replace(' ', '_')}_{uuid.uuid4().hex[:6]}.docx"
    out_path = _REPORT_OUTPUT_DIR / fname

    # ── TEMPLATE-FIRST PATH ──────────────────────────────────────────────
    # If a user-authored Word template exists at config/templates/office/cac/,
    # render that with {{placeholders}} → produces a polished docx that
    # inherits the user's brand, images, layout. Otherwise fall back to the
    # deterministic programmatic builder (kept as backup).
    try:
        pack = pack_from_workbook(workbook)
    except Exception as exc:  # noqa: BLE001
        log.error("cac_meeting_report.pack_failed", error=str(exc))
        raise HTTPException(500, f"Could not parse the data pack: {exc}")

    template_used = None
    try:
        from services.shared.office_template import (
            template_path_for, render_docx,
        )
        from services.shared.cac_report_docx import cac_report_context

        tpl = template_path_for("cac", "report", "CAC_Monthly_Report") \
            or template_path_for("cac", "report", "monthly")
        if tpl is not None:
            context, breaches = cac_report_context(pack, month)
            result = render_docx(tpl, context, out_path)
            template_used = str(tpl)
            log.info("cac_meeting_report.template_rendered",
                     template=template_used,
                     substitutions=result["substitutions"],
                     missing=result["missing"])
        else:
            # Fallback: existing deterministic builder
            breaches = build_cac_report_docx(pack, out_path, month)
            log.info("cac_meeting_report.programmatic_built",
                     reason="no user template at config/templates/office/cac/")
    except Exception as exc:  # noqa: BLE001
        log.error("cac_meeting_report.build_failed", error=str(exc),
                  template_used=template_used)
        raise HTTPException(500, f"Could not build the CAC report: {exc}")

    hard = [b for b in breaches if b.get("status") == "BREACH"]
    summary = (f"{len(hard)} limit breach(es), {len(breaches) - len(hard)} watch item(s)"
               if breaches else "no limit breaches")
    log.info("cac_meeting_report.done", month=month, file=fname, breaches=len(breaches))
    return {
        "answer": (f":page_facing_up: *CAC Meeting Report — First Draft ({month})* is ready "
                   f"(DRAFT for committee review) — {summary}. Figures are read straight from "
                   "the live Data Pack; see the attached .docx."),
        "confidence": "High",
        "file_path": str(out_path),
        "file_name": fname,
        "file_url": f"http://deck-writer:3050/reports/{fname}",
        "breaches": [f"{b['metric']} = {b['value']} (limit {b['limit']})" for b in breaches],
        "month": month,
    }
