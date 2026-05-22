# services/rag-ingestion/src/chunker.py
"""Document chunking pipeline for PDF, DOCX, XLSX, Markdown, and plain text."""
from __future__ import annotations

import hashlib
from datetime import datetime, timezone
from pathlib import Path

import structlog

from .config import RAGSettings

try:
    from services.shared.table_chunker import chunk_document as table_aware_chunk
    from services.shared.table_chunker import chunk_excel_sheet
except ImportError:
    table_aware_chunk = None
    chunk_excel_sheet = None

logger = structlog.get_logger("rag-ingestion.chunker")


class TextChunk:
    """A chunk of text with metadata."""

    __slots__ = ("text", "metadata")

    def __init__(self, text: str, metadata: dict[str, str | int | None]) -> None:
        self.text = text
        self.metadata = metadata


class DocumentChunker:
    """Chunks documents into text segments with metadata."""

    def __init__(self, settings: RAGSettings) -> None:
        self._chunk_size = settings.chunk_size
        self._overlap = settings.chunk_overlap

    async def chunk_file(
        self,
        file_path: Path,
        doc_type: str,
        dept: str = "CAC",
        extra_meta: dict[str, str] | None = None,
    ) -> list[TextChunk]:
        """Chunk a file into TextChunks with metadata."""
        handler = self._get_handler(doc_type)
        if handler is None:
            logger.warning("chunker.unsupported_type", doc_type=doc_type, path=str(file_path))
            return []

        try:
            raw_sections = handler(file_path)
        except Exception as exc:
            logger.error("chunker.extract_failed", path=str(file_path), error=str(exc))
            return []

        if not raw_sections:
            return []

        file_hash = self._hash_file(file_path)
        file_ext = file_path.suffix.lower()
        # ingested_at = wall-clock when this chunker ran. ISO-8601 UTC.
        # Lets retrieve_context filter by recency ("what changed this week").
        ingested_at = datetime.now(timezone.utc).isoformat()
        chunks: list[TextChunk] = []
        for section in raw_sections:
            text = section["text"].strip()
            if not text:
                continue

            # Use table-aware chunking if available and document contains tables
            if table_aware_chunk is not None and ("|---|" in text or file_ext in (".xlsx", ".csv")):
                table_chunks = table_aware_chunk(text, source=str(file_path), max_chunk_size=self._chunk_size)
                for c in table_chunks:
                    meta: dict[str, str | int | None] = {
                        "source_file": str(file_path),
                        "file_hash": file_hash,
                        "doc_type": doc_type,
                        "dept": dept,
                        "page": section.get("page"),
                        "section": section.get("section"),
                        "sheet": section.get("sheet"),
                        "chunk_type": c.chunk_type,
                        "ingested_at": ingested_at,
                    }
                    if extra_meta:
                        for k, v in extra_meta.items():
                            if v:
                                meta[k] = v
                    if c.metadata:
                        for k, v in c.metadata.items():
                            meta[k] = v
                    chunks.append(TextChunk(text=c.text, metadata=meta))
                continue

            for piece in self._split_text(text):
                meta = {
                    "source_file": str(file_path),
                    "file_hash": file_hash,
                    "doc_type": doc_type,
                    "dept": dept,
                    "page": section.get("page"),
                    "section": section.get("section"),
                    "sheet": section.get("sheet"),
                    "ingested_at": ingested_at,
                }
                if extra_meta:
                    for k, v in extra_meta.items():
                        if v:  # only add non-empty values
                            meta[k] = v
                chunks.append(TextChunk(text=piece, metadata=meta))

        logger.info("chunker.done", path=str(file_path), chunks=len(chunks))
        return chunks

    def _get_handler(self, doc_type: str):  # noqa: ANN202
        handlers = {
            "pdf": self._extract_pdf,
            "docx": self._extract_docx,
            "xlsx": self._extract_xlsx,
            "pptx": self._extract_pptx,
            "doc":  self._extract_doc,
            "xls":  self._extract_xls,
            "msg":  self._extract_msg,
            "md":   self._extract_text,
            "txt":  self._extract_text,
        }
        return handlers.get(doc_type)

    def _extract_pdf(self, path: Path) -> list[dict]:
        """Extract text from a PDF. If the file has no text layer (i.e. it's a
        scanned image PDF), fall back to OCR via ocrmypdf + tesseract before
        re-extracting. Tesseract languages are configured via the OCR_LANGS
        env var (default 'eng+tha' to handle the Thai SEC corpus)."""
        import fitz  # PyMuPDF
        import os

        def _read(p: str) -> list[dict]:
            doc = fitz.open(p)
            out: list[dict] = []
            try:
                for i, page in enumerate(doc):
                    text = page.get_text()
                    if text.strip():
                        out.append({"text": text, "page": i + 1, "section": None})
            finally:
                doc.close()
            return out

        sections = _read(str(path))
        if sections:
            return sections

        # No text layer — OCR fallback.
        try:
            import ocrmypdf
        except ImportError:
            logger.warning("chunker.ocr_unavailable", path=str(path),
                            msg="ocrmypdf not installed; cannot recover scanned PDF")
            return []

        langs = os.getenv("OCR_LANGS", "eng+tha")
        max_ocr_pages = int(os.getenv("OCR_MAX_PAGES", "60"))
        ocr_timeout = int(os.getenv("OCR_TIMEOUT_SECONDS", "300"))

        # Bound work: skip enormous PDFs to avoid runaway OCR jobs.
        try:
            page_count = fitz.open(str(path)).page_count
        except Exception:
            page_count = 0
        if page_count > max_ocr_pages:
            logger.info("chunker.ocr_skipped_too_large",
                         path=str(path), pages=page_count, max=max_ocr_pages)
            return []

        import tempfile, subprocess, time
        start = time.monotonic()
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            ocr_out = tmp.name
        try:
            # Run ocrmypdf as a subprocess so we can enforce a hard timeout.
            cmd = [
                "ocrmypdf", "--language", langs,
                "--output-type", "pdf",
                "--skip-text",   # don't re-OCR pages that already have text
                "--force-ocr",   # ...except the whole doc, since we got 0 text
                "--quiet",
                str(path), ocr_out,
            ]
            # --skip-text + --force-ocr are mutually exclusive; pick --force-ocr.
            cmd = [c for c in cmd if c != "--skip-text"]
            try:
                subprocess.run(cmd, check=True, timeout=ocr_timeout,
                                capture_output=True)
            except subprocess.TimeoutExpired:
                logger.warning("chunker.ocr_timeout",
                                path=str(path), seconds=ocr_timeout)
                return []
            except subprocess.CalledProcessError as exc:
                logger.warning("chunker.ocr_failed",
                                path=str(path),
                                stderr=(exc.stderr or b"").decode("utf-8", errors="replace")[:300])
                return []
            elapsed = time.monotonic() - start
            recovered = _read(ocr_out)
            logger.info("chunker.ocr_done",
                         path=str(path), elapsed_s=round(elapsed, 1),
                         sections=len(recovered), pages=page_count)
            return recovered
        finally:
            try:
                Path(ocr_out).unlink(missing_ok=True)
            except Exception:
                pass

    def _extract_docx(self, path: Path) -> list[dict]:
        from docx import Document

        doc = Document(str(path))
        sections = []
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                sections.append({"text": para.text, "page": None, "section": f"para_{i}"})
        return sections

    def _extract_xlsx(self, path: Path) -> list[dict]:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
        sections = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_vals = [c for c in row if c is not None]
                row_text = " | ".join(str(c) for c in row_vals)
                if row_text.strip():
                    rows.append(row_text)

            # Use table-aware Excel chunker if available
            if chunk_excel_sheet is not None and rows:
                structured_rows = []
                for row in ws.iter_rows(values_only=True):
                    structured_rows.append([str(c) if c is not None else "" for c in row])
                excel_chunks = chunk_excel_sheet(
                    structured_rows,
                    sheet_name=sheet_name,
                    source=str(path),
                    max_rows_per_chunk=self._chunk_size,
                )
                for c in excel_chunks:
                    sections.append({
                        "text": c.text,
                        "page": None,
                        "section": None,
                        "sheet": sheet_name,
                        "chunk_type": c.chunk_type,
                    })
            elif rows:
                sections.append({
                    "text": "\n".join(rows),
                    "page": None,
                    "section": None,
                    "sheet": sheet_name,
                })
        wb.close()
        return sections

    def _extract_text(self, path: Path) -> list[dict]:
        text = path.read_text(encoding="utf-8", errors="replace")
        return [{"text": text, "page": None, "section": None}] if text.strip() else []

    # --- PPTX (python-pptx) ---------------------------------------------
    def _extract_pptx(self, path: Path) -> list[dict]:
        from pptx import Presentation

        prs = Presentation(str(path))
        sections: list[dict] = []
        for i, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            # Slide title
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    parts.append(slide.shapes.title.text.strip())
            except Exception:
                pass
            # Body text from every shape with a text_frame
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            parts.append(t)
                # Tables
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
            # Speaker notes
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[notes] {notes}")
            except Exception:
                pass
            joined = "\n".join(parts).strip()
            if joined:
                sections.append({"text": joined, "page": i, "section": f"slide_{i}"})
        return sections

    # --- DOC (old binary Word, via antiword) ----------------------------
    def _extract_doc(self, path: Path) -> list[dict]:
        import subprocess
        try:
            r = subprocess.run(
                ["antiword", "-w", "0", str(path)],
                capture_output=True, timeout=60,
            )
            text = (r.stdout or b"").decode("utf-8", errors="replace")
        except FileNotFoundError as exc:
            logger.error("chunker.antiword_missing", path=str(path), error=str(exc))
            return []
        except subprocess.TimeoutExpired:
            logger.error("chunker.antiword_timeout", path=str(path))
            return []
        if not text.strip():
            return []
        # Split into paragraph-shaped sections so the downstream splitter has work.
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        return [
            {"text": p, "page": None, "section": f"para_{i}"}
            for i, p in enumerate(paragraphs)
        ] or [{"text": text, "page": None, "section": None}]

    # --- XLS (old binary Excel, via xlrd<2) -----------------------------
    def _extract_xls(self, path: Path) -> list[dict]:
        try:
            import xlrd
        except ImportError as exc:
            logger.error("chunker.xlrd_missing", error=str(exc))
            return []
        wb = xlrd.open_workbook(str(path))
        sections: list[dict] = []
        for sheet_name in wb.sheet_names():
            ws = wb.sheet_by_name(sheet_name)
            rows: list[str] = []
            for r in range(ws.nrows):
                vals = [str(ws.cell_value(r, c)).strip() for c in range(ws.ncols)]
                vals = [v for v in vals if v]
                if vals:
                    rows.append(" | ".join(vals))
            if rows:
                sections.append({
                    "text": "\n".join(rows),
                    "page": None,
                    "section": None,
                    "sheet": sheet_name,
                })
        return sections

    # --- MSG (Outlook email) --------------------------------------------
    def _extract_msg(self, path: Path) -> list[dict]:
        try:
            import extract_msg
        except ImportError as exc:
            logger.error("chunker.extract_msg_missing", error=str(exc))
            return []
        try:
            msg = extract_msg.Message(str(path))
        except Exception as exc:
            logger.error("chunker.msg_open_failed", path=str(path), error=str(exc))
            return []
        header_parts: list[str] = []
        if msg.sender:    header_parts.append(f"From: {msg.sender}")
        if msg.to:        header_parts.append(f"To: {msg.to}")
        if msg.cc:        header_parts.append(f"Cc: {msg.cc}")
        if msg.date:      header_parts.append(f"Date: {msg.date}")
        if msg.subject:   header_parts.append(f"Subject: {msg.subject}")
        body = (msg.body or "").strip()
        text = "\n".join(header_parts) + ("\n\n" + body if body else "")
        sections = [{"text": text, "page": None, "section": "message"}] if text.strip() else []
        # Attach inline list of attachment filenames as a separate section so it's searchable.
        try:
            attach_names = [a.longFilename or a.shortFilename or "unknown" for a in msg.attachments]
            if attach_names:
                sections.append({
                    "text": "Attachments: " + ", ".join(attach_names),
                    "page": None,
                    "section": "attachments",
                })
        except Exception:
            pass
        return sections

    def _split_text(self, text: str) -> list[str]:
        """Simple character-based splitter with overlap."""
        if len(text) <= self._chunk_size:
            return [text]
        pieces = []
        start = 0
        while start < len(text):
            end = start + self._chunk_size
            pieces.append(text[start:end])
            start = end - self._overlap
        return pieces

    @staticmethod
    def _hash_file(path: Path) -> str:
        h = hashlib.sha256()
        with open(path, "rb") as f:
            for block in iter(lambda: f.read(8192), b""):
                h.update(block)
        return h.hexdigest()
