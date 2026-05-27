"""Structured-table rendering into Word (.docx) and PowerPoint (.pptx).

Public API
----------
validate_table_spec(spec)        -> TableSpec
add_table_to_docx(doc, spec)     -> int
add_table_to_pptx(prs, spec)     -> int

TableSpec keys
--------------
title             (optional) bold caption rendered above the table
headers           required  list[str]
rows              required  list[list]; cell values coerced to str
column_widths_in  (optional) inches per column; must match headers length
style             "default" | "compact" | "emphasis"
align             (optional) per-column "left" | "right" | "center"

Docx header row: bold + light-grey fill (#E7E6E6).
Docx body rows:  alternating subtle banding (#F2F2F2 every other row).
Pptx header row: Brooker navy #0F3D5C background, white bold text.
"""
from __future__ import annotations

import logging
from typing import TypedDict

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# TypedDict contract
# ---------------------------------------------------------------------------

class TableSpec(TypedDict, total=False):
    title: str
    headers: list[str]
    rows: list[list]
    column_widths_in: list[float]
    style: str
    align: list[str]


# ---------------------------------------------------------------------------
# Validation
# ---------------------------------------------------------------------------

_VALID_STYLES = {"default", "compact", "emphasis"}
_VALID_ALIGNS = {"left", "right", "center"}


def validate_table_spec(spec: dict) -> TableSpec:
    """Validate a raw dict against TableSpec rules.

    Raises ValueError with a descriptive message on any violation.
    Coerces all cell values (headers and body cells) to str in-place on a copy.
    Returns the validated and coerced TableSpec.
    """
    if not isinstance(spec, dict):
        raise ValueError("TableSpec must be a dict")

    headers = spec.get("headers")
    if not headers or not isinstance(headers, list) or len(headers) == 0:
        raise ValueError("TableSpec.headers must be a non-empty list")

    rows = spec.get("rows")
    if rows is None or not isinstance(rows, list):
        raise ValueError("TableSpec.rows must be a list")

    n_cols = len(headers)

    # Coerce headers to str
    coerced_headers: list[str] = [str(h) for h in headers]

    # Validate and coerce rows
    coerced_rows: list[list[str]] = []
    for i, row in enumerate(rows):
        if not isinstance(row, list):
            raise ValueError(f"Row {i} must be a list, got {type(row).__name__}")
        if len(row) != n_cols:
            raise ValueError(
                f"Row {i} has {len(row)} cell(s) but headers define {n_cols} column(s)"
            )
        coerced_rows.append([str(cell) for cell in row])

    # Optional fields
    out: TableSpec = {
        "headers": coerced_headers,
        "rows": coerced_rows,
    }

    if "title" in spec:
        out["title"] = str(spec["title"])

    style = spec.get("style", "default")
    if style not in _VALID_STYLES:
        raise ValueError(
            f"TableSpec.style must be one of {_VALID_STYLES}, got {style!r}"
        )
    out["style"] = style

    col_widths = spec.get("column_widths_in")
    if col_widths is not None:
        if not isinstance(col_widths, list) or len(col_widths) != n_cols:
            raise ValueError(
                f"TableSpec.column_widths_in must be a list of {n_cols} floats "
                f"matching headers length"
            )
        out["column_widths_in"] = [float(w) for w in col_widths]

    align = spec.get("align")
    if align is not None:
        if not isinstance(align, list) or len(align) != n_cols:
            raise ValueError(
                f"TableSpec.align must be a list of {n_cols} alignment strings"
            )
        for a in align:
            if a not in _VALID_ALIGNS:
                raise ValueError(
                    f"TableSpec.align values must be one of {_VALID_ALIGNS}, got {a!r}"
                )
        out["align"] = list(align)

    return out


# ---------------------------------------------------------------------------
# Docx helpers
# ---------------------------------------------------------------------------

_HEADER_FILL = "E7E6E6"          # light grey
_BAND_FILL = "F2F2F2"            # subtle alternating row band
_BROOKER_NAVY = "0F3D5C"


def _set_cell_shading(cell, hex_color: str) -> None:
    """Apply a solid background fill to a python-docx table cell via lxml."""
    from lxml import etree

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove any existing shading element
    for existing in tcPr.findall(
        "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}shd"
    ):
        tcPr.remove(existing)

    shd = etree.SubElement(
        tcPr,
        "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}shd",
    )
    shd.set(
        "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val", "clear"
    )
    shd.set(
        "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}color", "auto"
    )
    shd.set(
        "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill",
        hex_color,
    )


def _set_cell_alignment(cell, align_str: str) -> None:
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    mapping = {
        "left": WD_ALIGN_PARAGRAPH.LEFT,
        "center": WD_ALIGN_PARAGRAPH.CENTER,
        "right": WD_ALIGN_PARAGRAPH.RIGHT,
    }
    align_val = mapping.get(align_str, WD_ALIGN_PARAGRAPH.LEFT)
    for para in cell.paragraphs:
        para.alignment = align_val


def _set_column_widths_docx(table, widths_in: list[float]) -> None:
    """Set per-cell column widths (python-docx requires per-cell, not per-column)."""
    from docx.shared import Inches

    for row in table.rows:
        for i, cell in enumerate(row.cells):
            if i < len(widths_in):
                cell.width = Inches(widths_in[i])


def _find_paragraph_index_by_heading(doc, hint: str) -> int | None:
    """Return the index (in doc.paragraphs) of the first heading containing hint."""
    for i, para in enumerate(doc.paragraphs):
        style_name = para.style.name.lower() if para.style else ""
        if "heading" in style_name and hint.lower() in para.text.lower():
            return i
    return None


def _insert_table_after_paragraph_index(doc, table, para_index: int) -> None:
    """Move `table` so that it appears immediately after doc.paragraphs[para_index].

    python-docx add_table() always appends to the body. We then detach the table's
    XML element and re-insert it after the target paragraph element.
    """
    from lxml import etree  # noqa: F401 (used via element ops)

    body = doc.element.body
    tbl_element = table._tbl

    # Detach from current position (it was appended to body)
    body.remove(tbl_element)

    # Find the target paragraph element
    target_para_element = doc.paragraphs[para_index]._p

    # Insert immediately after
    target_para_element.addnext(tbl_element)


# ---------------------------------------------------------------------------
# Public: add_table_to_docx
# ---------------------------------------------------------------------------

def add_table_to_docx(
    doc,
    spec: TableSpec,
    *,
    section_hint: str | None = None,
) -> int:
    """Insert a styled table into a python-docx Document.

    If section_hint is provided, substring-matches heading text and inserts the
    table (plus optional title paragraph) immediately after that heading.
    Otherwise appends at end.

    Header row: bold + light-grey fill (#E7E6E6).
    Body rows: alternating banding (#F2F2F2 on even rows, 0-indexed).

    Returns 1 on success. Raises ValueError on invalid spec.
    """
    validated = validate_table_spec(dict(spec))

    headers: list[str] = validated["headers"]
    rows: list[list[str]] = validated["rows"]
    n_cols = len(headers)
    n_rows = len(rows)
    style: str = validated.get("style", "default")
    col_widths = validated.get("column_widths_in")
    align = validated.get("align")
    title = validated.get("title")

    # Find insertion point
    insert_after_idx: int | None = None
    if section_hint is not None:
        insert_after_idx = _find_paragraph_index_by_heading(doc, section_hint)
        if insert_after_idx is None:
            logger.warning(
                "add_table_to_docx: section_hint %r not found; appending at end",
                section_hint,
            )

    # Optional caption paragraph — created before the table so we can move it too
    caption_para = None
    if title:
        caption_para = doc.add_paragraph()
        caption_para.add_run(title).bold = True

    # Create the table (always appended to body initially)
    tbl = doc.add_table(rows=1 + n_rows, cols=n_cols)

    # Apply named style if available; fall back gracefully
    try:
        tbl.style = "Light Grid Accent 1"
    except Exception:
        logger.debug("Table style 'Light Grid Accent 1' not available; using manual shading")

    # Header row
    hdr_cells = tbl.rows[0].cells
    for i, h in enumerate(headers):
        cell = hdr_cells[i]
        cell.text = ""
        run = cell.paragraphs[0].add_run(h)
        run.bold = True
        _set_cell_shading(cell, _HEADER_FILL)
        if align:
            _set_cell_alignment(cell, align[i])

    # Body rows
    for r_idx, row_data in enumerate(rows):
        tbl_row = tbl.rows[1 + r_idx]
        for c_idx, value in enumerate(row_data):
            cell = tbl_row.cells[c_idx]
            cell.text = str(value)
            if r_idx % 2 == 1:  # alternate banding on odd rows
                _set_cell_shading(cell, _BAND_FILL)
            if align:
                _set_cell_alignment(cell, align[c_idx])

    # Column widths (must be set per-cell after all rows are populated)
    if col_widths:
        _set_column_widths_docx(tbl, col_widths)

    # Relocate to after the target heading if section_hint matched
    if insert_after_idx is not None:
        if caption_para is not None:
            # Move caption first, then the table after the caption
            _insert_table_after_paragraph_index(doc, tbl, insert_after_idx)
            # caption_para is a paragraph already in the body; move it before tbl
            body = doc.element.body
            cap_p = caption_para._p
            tbl_element = tbl._tbl
            body.remove(cap_p)
            tbl_element.addprevious(cap_p)
        else:
            _insert_table_after_paragraph_index(doc, tbl, insert_after_idx)

    return 1


# ---------------------------------------------------------------------------
# Pptx helpers
# ---------------------------------------------------------------------------

def _find_slide_by_title_hint(prs, hint: str) -> int | None:
    """Return the slide index whose title text contains hint (case-insensitive)."""
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and hint.lower() in shape.text_frame.text.lower():
                return i
    return None


def _add_blank_slide(prs):
    """Add a blank slide using the blank layout (last layout, or index 6 fallback)."""
    slide_layouts = prs.slide_layouts
    # Try to find a blank layout
    blank_layout = None
    for layout in slide_layouts:
        if layout.name.lower() in ("blank", "blank slide"):
            blank_layout = layout
            break
    if blank_layout is None:
        # Fall back to the last layout which is typically blank
        blank_layout = slide_layouts[-1]
    return prs.slides.add_slide(blank_layout)


# ---------------------------------------------------------------------------
# Public: add_table_to_pptx
# ---------------------------------------------------------------------------

def add_table_to_pptx(
    prs,
    spec: TableSpec,
    *,
    slide_index: int | None = None,
    slide_title_hint: str | None = None,
    left_inches: float = 0.5,
    top_inches: float = 1.5,
    width_inches: float = 9.0,
    height_inches: float = 4.5,
) -> int:
    """Insert a table shape into a python-pptx Presentation slide.

    Priority: slide_index > slide_title_hint > new slide.
    If neither resolves, a new blank slide is created.

    Title row: Brooker navy #0F3D5C background, white bold text.
    Body rows: alternating subtle banding.

    Returns 1 on success. Raises ValueError on invalid spec.
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    validated = validate_table_spec(dict(spec))

    headers: list[str] = validated["headers"]
    rows: list[list[str]] = validated["rows"]
    n_cols = len(headers)
    n_rows = len(rows)
    align = validated.get("align")

    # Resolve target slide
    target_slide = None

    if slide_index is not None:
        if 0 <= slide_index < len(prs.slides):
            target_slide = prs.slides[slide_index]
        else:
            logger.warning(
                "add_table_to_pptx: slide_index %d out of range (%d slides); "
                "creating new slide",
                slide_index,
                len(prs.slides),
            )

    if target_slide is None and slide_title_hint is not None:
        idx = _find_slide_by_title_hint(prs, slide_title_hint)
        if idx is not None:
            target_slide = prs.slides[idx]
        else:
            logger.warning(
                "add_table_to_pptx: slide_title_hint %r not found; creating new slide",
                slide_title_hint,
            )

    if target_slide is None:
        target_slide = _add_blank_slide(prs)

    # Add table shape
    left = Inches(left_inches)
    top = Inches(top_inches)
    width = Inches(width_inches)
    height = Inches(height_inches)

    table_shape = target_slide.shapes.add_table(
        1 + n_rows, n_cols, left, top, width, height
    )
    table = table_shape.table

    # Resolve per-column widths proportionally if column_widths_in provided
    col_widths = validated.get("column_widths_in")
    if col_widths:
        total_spec = sum(col_widths)
        total_emu = Inches(width_inches)
        for c_idx in range(n_cols):
            table.columns[c_idx].width = int(
                total_emu * col_widths[c_idx] / total_spec
            )

    # Header row — navy background, white bold text
    navy = RGBColor(0x0F, 0x3D, 0x5C)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    band_color = RGBColor(0xF2, 0xF2, 0xF2)

    for c_idx, header_text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = navy

        tf = cell.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.color.rgb = white
            if align:
                _apply_pptx_para_align(para, align[c_idx])

    # Body rows
    for r_idx, row_data in enumerate(rows):
        for c_idx, value in enumerate(row_data):
            cell = table.cell(1 + r_idx, c_idx)
            cell.text = str(value)

            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = band_color

            if align:
                tf = cell.text_frame
                for para in tf.paragraphs:
                    _apply_pptx_para_align(para, align[c_idx])

    return 1


def _apply_pptx_para_align(para, align_str: str) -> None:
    from pptx.enum.text import PP_ALIGN

    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    val = mapping.get(align_str)
    if val is not None:
        para.alignment = val
