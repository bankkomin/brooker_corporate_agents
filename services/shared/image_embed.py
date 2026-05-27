"""Image embedding helpers for Word (.docx) and PowerPoint (.pptx) documents.

Resolves an image from three possible sources (HTTP URL, filesystem path, or
MinIO object key), validates the bytes are a real PNG or JPEG, then inserts the
image at the right place in the document.

Follows the same idioms as office_template.py and cac_report_docx.py:
- Local imports for heavy deps so callers that don't need them stay light.
- Fail-fast with descriptive errors; never swallow exceptions silently.
- No hardcoded credentials — MinIO config comes from env vars.

Public API
----------
    fetch_image_bytes(src)            -> bytes
    embed_image_in_docx(doc, src, ...) -> int  (1 = inserted, 0 = skipped)
    embed_image_in_pptx(prs, src, ...) -> int  (always 1)
"""
from __future__ import annotations

import io
import os
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    pass  # heavy imports (docx, pptx) happen inside functions

try:
    import structlog
    log = structlog.get_logger(__name__)
except ImportError:
    import logging
    log = logging.getLogger(__name__)

# ImageSource is a plain dict with exactly one of these keys.
ImageSource = dict  # {"url": str} | {"path": str} | {"minio_key": str}

# ----- constants -------------------------------------------------------------

_MAX_BYTES: int = 25 * 1024 * 1024          # 25 MB hard cap
_HTTP_TIMEOUT: float = 10.0                  # seconds
_MINIO_BUCKET: str = "paperclip-uploads"

# Valid magic-byte prefixes.  We intentionally keep this small — the test spec
# requires PNG and JPEG; GIF / WebP etc. are rejected with a clear message.
_MAGIC: dict[bytes, str] = {
    b"\x89PNG": "PNG",
    b"\xff\xd8": "JPEG",
}


# ----- image resolution ------------------------------------------------------

def _check_magic(data: bytes) -> None:
    """Raise ValueError if data doesn't look like a PNG or JPEG."""
    for prefix, fmt in _MAGIC.items():
        if data[: len(prefix)] == prefix:
            log.debug("image.magic_ok", format=fmt, size=len(data))
            return
    first4 = data[:4].hex()
    raise ValueError(
        f"Unsupported image format (magic bytes: {first4!r}). "
        "Only PNG and JPEG are accepted."
    )


def fetch_image_bytes(src: ImageSource) -> bytes:
    """Resolve *src* to raw PNG/JPEG bytes.

    Supported source shapes
    -----------------------
    ``{"url": "https://..."}``
        Fetched via ``httpx`` with a 10-second timeout.  Content-Length (when
        present) is checked *before* the download; the response body is also
        capped at 25 MB during streaming so oversized files never fully buffer.

    ``{"path": "/some/file.png"}``
        Read straight from the filesystem.  Both absolute and relative paths
        are accepted; a missing file raises ``ValueError``.

    ``{"minio_key": "uploads/abc.png"}``
        Downloaded from the ``paperclip-uploads`` MinIO bucket.  Endpoint and
        credentials come from env vars:
        ``MINIO_ENDPOINT`` (default ``minio:9000``),
        ``MINIO_ACCESS_KEY``, ``MINIO_SECRET_KEY``.

    Raises
    ------
    ValueError
        Bad/missing src key, oversized content-length, bad magic bytes, or
        file-not-found for path-based sources.
    """
    if not isinstance(src, dict) or not src:
        raise ValueError("src must be a non-empty dict with one of: url, path, minio_key")

    if "url" in src:
        return _fetch_from_url(str(src["url"]))
    if "path" in src:
        return _fetch_from_path(str(src["path"]))
    if "minio_key" in src:
        return _fetch_from_minio(str(src["minio_key"]))

    raise ValueError(
        f"Unrecognised ImageSource key(s): {list(src.keys())}. "
        "Expected exactly one of: url, path, minio_key"
    )


def _fetch_from_url(url: str) -> bytes:
    import httpx

    if not url.lower().startswith(("http://", "https://")):
        raise ValueError(f"URL must start with http:// or https://, got: {url!r}")

    log.info("image.fetch_url", url=url[:120])
    with httpx.Client(timeout=_HTTP_TIMEOUT) as client:
        # HEAD first to check Content-Length without downloading the body.
        # Not all servers honour HEAD, so we only gate on it when the header is
        # actually present; the streaming check below is the definitive guard.
        try:
            head = client.head(url, follow_redirects=True)
            cl = int(head.headers.get("content-length", 0))
            if cl > _MAX_BYTES:
                raise ValueError(
                    f"Remote image is too large: Content-Length {cl} bytes "
                    f"exceeds the 25 MB cap ({_MAX_BYTES} bytes)."
                )
        except httpx.HTTPStatusError:
            pass  # HEAD rejected — proceed to GET
        except ValueError:
            raise  # re-raise our own size error
        except Exception:
            pass  # HEAD failed for other reason — proceed to GET

        # Stream the GET response so we can bail out mid-download if needed.
        with client.stream("GET", url, follow_redirects=True) as resp:
            resp.raise_for_status()
            cl_str = resp.headers.get("content-length", "")
            if cl_str:
                try:
                    cl = int(cl_str)
                    if cl > _MAX_BYTES:
                        raise ValueError(
                            f"Remote image is too large: Content-Length {cl} bytes "
                            f"exceeds the 25 MB cap ({_MAX_BYTES} bytes)."
                        )
                except ValueError as exc:
                    if "too large" in str(exc):
                        raise

            chunks: list[bytes] = []
            total = 0
            for chunk in resp.iter_bytes(chunk_size=65536):
                total += len(chunk)
                if total > _MAX_BYTES:
                    raise ValueError(
                        f"Remote image body exceeded the 25 MB cap while streaming "
                        f"(read {total} bytes so far)."
                    )
                chunks.append(chunk)

    data = b"".join(chunks)
    _check_magic(data)
    return data


def _fetch_from_path(path_str: str) -> bytes:
    from pathlib import Path

    p = Path(path_str)
    if not p.exists():
        raise ValueError(f"Image file not found: {path_str!r}")
    if not p.is_file():
        raise ValueError(f"Image path is not a file: {path_str!r}")

    log.info("image.fetch_path", path=path_str)
    data = p.read_bytes()
    if len(data) > _MAX_BYTES:
        raise ValueError(
            f"Image file is too large: {len(data)} bytes exceeds the 25 MB cap."
        )
    _check_magic(data)
    return data


def _fetch_from_minio(key: str) -> bytes:
    from minio import Minio
    from minio.error import S3Error

    endpoint = os.getenv("MINIO_ENDPOINT", "minio:9000")
    access_key = os.getenv("MINIO_ACCESS_KEY", "")
    secret_key = os.getenv("MINIO_SECRET_KEY", "")

    log.info("image.fetch_minio", key=key, bucket=_MINIO_BUCKET, endpoint=endpoint)
    client = Minio(endpoint, access_key=access_key, secret_key=secret_key, secure=False)

    try:
        response = client.get_object(_MINIO_BUCKET, key)
        try:
            data = response.read()
        finally:
            response.close()
            response.release_conn()
    except S3Error as exc:
        raise ValueError(f"MinIO fetch failed for key {key!r}: {exc}") from exc

    if len(data) > _MAX_BYTES:
        raise ValueError(
            f"MinIO object is too large: {len(data)} bytes exceeds the 25 MB cap."
        )
    _check_magic(data)
    return data


# ----- docx embedding --------------------------------------------------------

def embed_image_in_docx(
    doc,
    src: ImageSource,
    *,
    section_hint: str | None = None,
    caption: str | None = None,
    width_inches: float = 5.5,
) -> int:
    """Insert an image into a python-docx ``Document``.

    Parameters
    ----------
    doc
        A ``docx.Document`` instance (mutated in place).
    src
        ImageSource dict — passed straight to ``fetch_image_bytes``.
    section_hint
        Case-insensitive substring.  If any heading paragraph's text contains
        this string the image is inserted immediately *after* that heading.
        When ``None`` or no heading matches, the image is appended at the end.
    caption
        When given, a centred italic paragraph is added below the image.
    width_inches
        Rendered width of the image.  Height is scaled proportionally.

    Returns
    -------
    int
        ``1`` on success, ``0`` if the image was skipped (should not happen in
        normal usage — hard failures raise instead).
    """
    from docx.enum.text import WD_ALIGN_PARAGRAPH  # type: ignore[attr-defined]
    from docx.oxml.ns import qn
    from docx.shared import Inches
    from lxml import etree

    data = fetch_image_bytes(src)
    img_stream = io.BytesIO(data)

    # Find the insertion point.  We walk doc.paragraphs (body only; excludes
    # headers/footers) looking for a heading whose text contains section_hint.
    insert_after_paragraph = None
    if section_hint:
        hint_lower = section_hint.lower()
        for para in doc.paragraphs:
            if hint_lower in para.text.lower():
                insert_after_paragraph = para
                break

    if insert_after_paragraph is not None:
        # Insert a new paragraph directly after the matched heading using lxml.
        # python-docx doesn't expose an insert_paragraph_after helper, so we
        # operate on the underlying XML element.
        _insert_image_after_paragraph(
            doc, insert_after_paragraph, img_stream, width_inches,
            caption=caption,
        )
    else:
        # Append at the end of the document body.
        img_para = doc.add_paragraph()
        run = img_para.add_run()
        run.add_picture(img_stream, width=Inches(width_inches))
        _center_paragraph(img_para)
        if caption:
            _add_caption_paragraph(doc.add_paragraph(), caption)

    log.info("image.embedded_docx",
             section_hint=section_hint, caption=bool(caption),
             width_inches=width_inches)
    return 1


def _center_paragraph(para) -> None:
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    para.alignment = WD_ALIGN_PARAGRAPH.CENTER


def _add_caption_paragraph(para, caption: str) -> None:
    """Mutate *para* to be a centred italic caption."""
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    run = para.add_run(caption)
    run.italic = True
    para.alignment = WD_ALIGN_PARAGRAPH.CENTER


def _insert_image_after_paragraph(
    doc,
    anchor_para,
    img_stream: io.BytesIO,
    width_inches: float,
    *,
    caption: str | None,
) -> None:
    """Use lxml addnext to splice an image paragraph after *anchor_para*.

    python-docx's OxmlElement + _p / _Body approach:
      - Each paragraph is a ``<w:p>`` element.
      - The document body is the parent.
      - ``anchor_para._p.addnext(new_p._p)`` inserts after anchor in XML order.

    We fabricate the image paragraph by appending it (which goes to the end of
    the document), then relocating its ``<w:p>`` element via ``addnext``.  This
    way python-docx fully owns the picture relationship — we never have to touch
    the ZIP internals ourselves.
    """
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches

    # Step 1: append a temporary paragraph at the end — python-docx wires up
    # the image relationship cleanly.
    img_para = doc.add_paragraph()
    run = img_para.add_run()
    run.add_picture(img_stream, width=Inches(width_inches))
    _center_paragraph(img_para)

    # Step 2: physically move the <w:p> element to right after the anchor.
    anchor_para._p.addnext(img_para._p)

    # Step 3: caption paragraph (if requested) goes right after the image para.
    if caption:
        cap_para = doc.add_paragraph()
        _add_caption_paragraph(cap_para, caption)
        # Caption lands after the (now-repositioned) image para.
        img_para._p.addnext(cap_para._p)


# ----- pptx embedding --------------------------------------------------------

def embed_image_in_pptx(
    prs,
    src: ImageSource,
    *,
    slide_index: int | None = None,
    slide_title_hint: str | None = None,
    caption: str | None = None,
    left_inches: float = 0.5,
    top_inches: float = 1.5,
    width_inches: float = 6.0,
) -> int:
    """Insert an image into a python-pptx ``Presentation``.

    Target slide resolution order
    ------------------------------
    1. ``slide_index`` (0-based) — used directly if within range.
    2. ``slide_title_hint`` — case-insensitive substring match against each
       slide's title shape text.
    3. Neither provided / no match — a new blank slide is appended using
       ``prs.slide_layouts[5]`` (Title Only) and the image is placed there.

    Parameters
    ----------
    caption
        When given, a small text box is added below the image.
    left_inches, top_inches, width_inches
        Position and size of the image on the slide.

    Returns
    -------
    int
        Always ``1``.
    """
    from pptx.util import Inches, Pt

    data = fetch_image_bytes(src)
    img_stream = io.BytesIO(data)

    slide = _resolve_pptx_slide(prs, slide_index=slide_index,
                                  slide_title_hint=slide_title_hint)

    pic = slide.shapes.add_picture(
        img_stream,
        Inches(left_inches),
        Inches(top_inches),
        width=Inches(width_inches),
    )

    if caption:
        pic_bottom = top_inches + pic.height / 914400  # EMU -> inches
        txb = slide.shapes.add_textbox(
            Inches(left_inches),
            Inches(pic_bottom + 0.05),
            Inches(width_inches),
            Inches(0.4),
        )
        tf = txb.text_frame
        tf.text = caption
        for para in tf.paragraphs:
            from pptx.enum.text import PP_ALIGN
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.italic = True
                run.font.size = Pt(11)

    log.info("image.embedded_pptx",
             slide_index=slide_index, slide_title_hint=slide_title_hint,
             caption=bool(caption), width_inches=width_inches)
    return 1


def _resolve_pptx_slide(prs, *, slide_index: int | None, slide_title_hint: str | None):
    """Return the target slide, creating a new one if nothing matches."""
    slides = prs.slides

    # 1. Explicit index
    if slide_index is not None:
        if 0 <= slide_index < len(slides):
            return slides[slide_index]
        # Out-of-range: fall through to "new slide" path below.
        log.warning(
            "image.pptx_slide_index_out_of_range slide_index=%s total=%s",
            slide_index,
            len(slides),
        )

    # 2. Title hint — substring match
    if slide_title_hint:
        hint_lower = slide_title_hint.lower()
        for slide in slides:
            title_shape = slide.shapes.title
            if title_shape and hint_lower in (title_shape.text or "").lower():
                return slide

    # 3. No match — append a new blank slide
    layout = prs.slide_layouts[5]  # Title Only (standard Office layout 5)
    return prs.slides.add_slide(layout)
