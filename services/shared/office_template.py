"""User-authored Microsoft Office template renderer.

Lets the user create polished templates in real Word / Excel / PowerPoint
(with images, logos, brand themes, real styles, charts...) and the agents
just fill in placeholders. Preserves everything: images, headers/footers,
slide masters, named cells, conditional formatting, charts.

Placeholder convention: `{{snake_case_key}}` typed as plain text anywhere in
the template. Examples:
  Word:   "Total assets: {{total_assets}}"
  Excel:  any cell containing "{{net_worth}}"
  PowerPoint: any text frame / table cell

Usage:
    from services.shared.office_template import render_docx, render_xlsx, render_pptx

    render_docx(
        template_path="config/templates/office/cac/CAC_Monthly_Report.docx",
        context={
            "month": "May 2026",
            "total_assets": "THB 3,127.97M",
            "net_worth": "THB 2,435.11M",
            ...
        },
        out_path="data/reports/CAC_Report_May_2026.docx",
    )

Notes:
- Placeholders that survive multi-run formatting work correctly (we join the
  paragraph text, do the substitution, then put the result back into the first
  run preserving its formatting).
- Missing placeholders are left as `{{key}}` in the output (visible defect by
  design — silent omission would let bad output through).
- Tables in Word + Excel + PowerPoint all get scanned.
- Headers/footers in Word are scanned.
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Any

# Single source of truth — placeholders look like {{snake_case_key}} with
# optional whitespace inside the braces. We deliberately reject keys with dots
# / brackets so a sentence containing "{{ie}} could match" doesn't false-fire.
_PLACEHOLDER_RE = re.compile(r"\{\{\s*([a-zA-Z_][a-zA-Z0-9_]*)\s*\}\}")


def _resolve(text: str, context: dict[str, Any]) -> tuple[str, int]:
    """Replace every {{key}} in `text` with str(context[key]).
    Returns (new_text, n_substitutions). Unknown keys are left as-is."""
    n = 0

    def sub(m: re.Match) -> str:
        nonlocal n
        key = m.group(1)
        if key in context:
            n += 1
            return str(context[key])
        return m.group(0)  # leave unknown placeholders visible

    return _PLACEHOLDER_RE.sub(sub, text), n


# ── Word (.docx) ──────────────────────────────────────────────────────────

def _docx_replace_in_paragraph(paragraph, context: dict[str, Any]) -> int:
    """Replace placeholders in a Word paragraph, even when split across runs.

    Strategy: concatenate all run text, do the substitution against the full
    paragraph, then put the new text into the first non-empty run and clear
    the others. Formatting of the FIRST run is preserved; mid-paragraph
    inline-formatting changes are flattened. For typical templates (placeholder
    in a single uniform-style sentence) this works fine.
    """
    full = "".join(r.text for r in paragraph.runs)
    new_text, n = _resolve(full, context)
    if n == 0 or new_text == full:
        return 0
    runs = paragraph.runs
    if runs:
        runs[0].text = new_text
        for r in runs[1:]:
            r.text = ""
    return n


def render_docx(template_path: str | Path, context: dict[str, Any],
                out_path: str | Path) -> dict[str, Any]:
    """Open a .docx template, substitute {{placeholders}}, save to out_path.

    Returns {"substitutions": N, "missing": [list of unresolved keys]}.
    """
    from docx import Document

    doc = Document(str(template_path))
    total = 0

    # Body paragraphs
    for p in doc.paragraphs:
        total += _docx_replace_in_paragraph(p, context)

    # Tables (recursive — tables can nest)
    def walk_table(table) -> int:
        n = 0
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    n += _docx_replace_in_paragraph(p, context)
                for inner in cell.tables:
                    n += walk_table(inner)
        return n

    for tbl in doc.tables:
        total += walk_table(tbl)

    # Headers + footers across all sections
    for section in doc.sections:
        for container in (section.header, section.footer,
                          section.first_page_header, section.first_page_footer,
                          section.even_page_header, section.even_page_footer):
            try:
                for p in container.paragraphs:
                    total += _docx_replace_in_paragraph(p, context)
                for tbl in container.tables:
                    total += walk_table(tbl)
            except Exception:
                pass  # not all section variants exist for every doc

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))

    # Find any unresolved placeholders by scanning the saved doc text
    full_text = "\n".join(p.text for p in doc.paragraphs)
    missing = sorted(set(_PLACEHOLDER_RE.findall(full_text)))
    return {"substitutions": total, "missing": missing,
            "template": str(template_path), "out": str(out_path)}


# ── Excel (.xlsx) ─────────────────────────────────────────────────────────

def render_xlsx(template_path: str | Path, context: dict[str, Any],
                out_path: str | Path) -> dict[str, Any]:
    """Open an .xlsx template, substitute {{placeholders}} in every cell.

    Also supports named ranges: if `context` has a key matching a named range
    in the workbook (e.g. context["TotalAssets"] = 3128 and the workbook has
    a named range "TotalAssets" pointing at G14), the cell is set directly.
    """
    import openpyxl

    wb = openpyxl.load_workbook(str(template_path))
    total = 0
    missing: set[str] = set()

    # Pass 1 — placeholder text in cells
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is None or not isinstance(cell.value, str):
                    continue
                new_v, n = _resolve(cell.value, context)
                if n:
                    cell.value = new_v
                    total += n
                else:
                    for m in _PLACEHOLDER_RE.finditer(cell.value):
                        missing.add(m.group(1))

    # Pass 2 — defined names (named ranges). Skips builtins like _xlnm.
    for name in wb.defined_names:
        if name.startswith("_xlnm"):
            continue
        if name not in context:
            continue
        try:
            dest = wb.defined_names[name].value  # e.g. "Sheet1!$G$14"
            for sheet_ref, cell_ref in wb.defined_names[name].destinations:
                ws = wb[sheet_ref]
                ws[cell_ref].value = context[name]
                total += 1
        except Exception:
            pass

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out_path))
    return {"substitutions": total, "missing": sorted(missing),
            "template": str(template_path), "out": str(out_path)}


# ── PowerPoint (.pptx) ────────────────────────────────────────────────────

def _pptx_replace_in_text_frame(text_frame, context: dict[str, Any]) -> int:
    n = 0
    for para in text_frame.paragraphs:
        full = "".join(r.text for r in para.runs)
        new_text, sub_n = _resolve(full, context)
        if sub_n and new_text != full:
            runs = para.runs
            if runs:
                runs[0].text = new_text
                for r in runs[1:]:
                    r.text = ""
            n += sub_n
    return n


def render_pptx(template_path: str | Path, context: dict[str, Any],
                out_path: str | Path) -> dict[str, Any]:
    """Open a .pptx template, substitute {{placeholders}} in every text frame
    and table cell across every slide. Slide masters / layouts / themes / images
    are preserved by python-pptx unchanged."""
    from pptx import Presentation

    prs = Presentation(str(template_path))
    total = 0
    missing: set[str] = set()

    for slide in prs.slides:
        for shape in slide.shapes:
            # Plain text frames
            if shape.has_text_frame:
                total += _pptx_replace_in_text_frame(shape.text_frame, context)
                for m in _PLACEHOLDER_RE.finditer(shape.text_frame.text):
                    missing.add(m.group(1))
            # Tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        total += _pptx_replace_in_text_frame(cell.text_frame, context)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return {"substitutions": total, "missing": sorted(missing),
            "template": str(template_path), "out": str(out_path)}


# ── Discovery helper ──────────────────────────────────────────────────────

def template_path_for(dept: str, kind: str, name: str,
                      base: str = "config/templates/office") -> Path | None:
    """Locate a user-authored template; return None if it doesn't exist.

    Caller pattern:
        tpl = template_path_for("cac", "report", "monthly")
        if tpl:
            render_docx(tpl, context, out)
        else:
            <existing programmatic builder as fallback>
    """
    ext = {"report": "docx", "minutes": "docx", "memo": "docx",
           "deck": "pptx", "presentation": "pptx",
           "tracker": "xlsx", "dashboard": "xlsx", "sheet": "xlsx"}.get(kind, kind)
    candidates = [
        Path(base) / dept / f"{name}.{ext}",
        Path(base) / dept / f"{name}_template.{ext}",
        Path(base) / dept / f"{name}-template.{ext}",
    ]
    for c in candidates:
        if c.is_file():
            return c
    return None
