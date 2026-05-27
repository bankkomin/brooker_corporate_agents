"""Integration tests: deck-writer rich-content endpoints.

Uses FastAPI's TestClient (no Docker rebuild required).  External services
(RAG ingestion, Qdrant, LLM) are mocked at the httpx.AsyncClient level so
the tests exercise endpoint plumbing end-to-end without network dependencies.

Tests:
    1. test_report_with_image_embed
       POST /report with one ImageEmbed pointing at a local PNG (created via
       PIL in the test); verifies the response is a .docx that contains at
       least one inline image shape.

    2. test_compose_with_chart_embed
       POST /compose with one ChartEmbed; verifies the response is a .pptx
       and has at least one picture shape across all slides.

    3. test_compose_xlsx_simple
       POST /compose-xlsx with 2 sheets and a formula; verifies the response
       is a .xlsx with both sheets present and the formula in the correct cell.
"""
from __future__ import annotations

import io
import json
import os
import tempfile
from pathlib import Path
from typing import Any
from unittest.mock import AsyncMock, MagicMock, patch

import pytest


# ---------------------------------------------------------------------------
# Module-level availability guards
# ---------------------------------------------------------------------------

def _require(*modules: str) -> None:
    """Skip the whole module if any required package is not importable."""
    missing = []
    for mod in modules:
        try:
            __import__(mod)
        except ImportError:
            missing.append(mod)
    if missing:
        pytest.skip(f"Required module(s) not importable: {', '.join(missing)}")


_require(
    "services.shared.image_embed",
    "services.shared.chart_render",
    "services.shared.table_render",
    "services.shared.xlsx_compose",
    "services.shared.mermaid_render",
    "services.shared.drafter_table_prompt",
    "fastapi",
    "pptx",
    "docx",
    "PIL",
    "openpyxl",
)


# ---------------------------------------------------------------------------
# Shared fixtures
# ---------------------------------------------------------------------------

_MINIMAL_OUTLINE = {
    "title": "Test Report",
    "subtitle": None,
    "outline": [
        {
            "narrative_role": "cover",
            "headline": "Test",
            "key_points": [],
            "supporting_data": [],
        },
        {
            "narrative_role": "context",
            "headline": "Background",
            "key_points": ["Point A", "Point B"],
            "supporting_data": [],
        },
        {
            "narrative_role": "conclusion",
            "headline": "Q&A",
            "key_points": [],
            "supporting_data": [],
        },
    ],
}

_MINIMAL_DECK_SPEC = {
    "title": "Test Deck",
    "subtitle": None,
    "slides": [
        {"layout": "cover", "title": "Test Deck", "bullets": []},
        {"layout": "bullets", "title": "Background", "bullets": ["Point A", "Point B"]},
        {"layout": "closing", "title": "Q&A", "bullets": []},
    ],
}


def _mock_http_responses() -> dict[str, Any]:
    """Return a mapping of URL-substring → mock response payload.

    The AsyncClient post() mock uses this to return canned responses for each
    downstream service call deck-writer makes during a request.
    """
    return {
        "/embed": {"vector": [0.1] * 768},
        "/collections": {"result": []},
        "/chat/completions": {
            "choices": [
                {
                    "message": {
                        "content": json.dumps(_MINIMAL_OUTLINE)
                    }
                }
            ]
        },
    }


def _make_async_client_mock(responses: dict[str, Any]) -> MagicMock:
    """Create a mock httpx.AsyncClient that returns canned responses."""
    mock_client = MagicMock()
    mock_client.__aenter__ = AsyncMock(return_value=mock_client)
    mock_client.__aexit__ = AsyncMock(return_value=None)
    mock_client.aclose = AsyncMock(return_value=None)

    async def _post(url: str, **kwargs: Any):
        resp = MagicMock()
        resp.raise_for_status = MagicMock()
        resp.status_code = 200
        for url_frag, payload in responses.items():
            if url_frag in url:
                resp.json = MagicMock(return_value=payload)
                return resp
        # Default: empty result
        resp.json = MagicMock(return_value={"result": [], "vector": [0.0] * 768})
        return resp

    mock_client.post = AsyncMock(side_effect=_post)
    return mock_client


def _make_minimal_png() -> bytes:
    """Return a minimal valid PNG as bytes using PIL."""
    from PIL import Image

    img = Image.new("RGB", (64, 64), color=(0, 120, 200))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Fixtures that set up the FastAPI TestClient with mocked env and HTTP client
# ---------------------------------------------------------------------------

@pytest.fixture()
def deck_writer_env(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    """Set env vars so deck-writer can start without Docker services."""
    monkeypatch.setenv("VLLM_LARGE_URL", "http://fake-llm:8000/v1")
    monkeypatch.setenv("VLLM_LARGE_MODEL", "test-model")
    monkeypatch.setenv("RAG_INGESTION_URL", "http://fake-rag:3004")
    monkeypatch.setenv("QDRANT_URL", "http://fake-qdrant:6333")
    monkeypatch.setenv("DECK_OUTPUT_DIR", str(tmp_path / "decks"))
    monkeypatch.setenv("LOG_LEVEL", "WARNING")
    (tmp_path / "decks").mkdir()
    (tmp_path / "reports").mkdir()
    # Point report output dir to tmp as well — patch the module-level constant
    # after import so the directory exists.
    return tmp_path


def _load_deck_writer_main():
    """Load services/deck-writer/src/main.py as a module.

    The directory name uses a hyphen which is illegal in Python identifiers,
    so standard `import services.deck_writer` fails.  We use importlib to
    load the module from its actual file path.

    The module is cached in sys.modules under the key ``deck_writer_main`` so
    it is executed only once per process.  This avoids Pydantic class-identity
    errors that arise when the same source file is exec'd multiple times,
    producing different class objects that Pydantic treats as incompatible.
    """
    import importlib.util
    import sys

    _MODULE_KEY = "deck_writer_main"
    if _MODULE_KEY in sys.modules:
        return sys.modules[_MODULE_KEY]

    here = Path(__file__).resolve().parents[2]  # project root
    spec = importlib.util.spec_from_file_location(
        _MODULE_KEY,
        here / "services" / "deck-writer" / "src" / "main.py",
    )
    if spec is None or spec.loader is None:
        raise ImportError("Could not locate services/deck-writer/src/main.py")
    mod = importlib.util.module_from_spec(spec)
    sys.modules[_MODULE_KEY] = mod
    spec.loader.exec_module(mod)  # type: ignore[union-attr]
    return mod


@pytest.fixture()
def client(deck_writer_env: Path):
    """Return a TestClient for the deck-writer FastAPI app.

    Mocking strategy: the lifespan creates a real httpx.AsyncClient and
    assigns it to ``main_mod._http``.  We intercept at a lower level by
    patching ``httpx.AsyncClient`` itself so that every instance returned
    during the test uses our canned responses.  This keeps the lifespan
    logic intact while preventing any real network calls.
    """
    from fastapi.testclient import TestClient

    main_mod = _load_deck_writer_main()

    # Redirect output dirs to tmp so tests don't need /data/decks or /data/reports.
    report_dir = deck_writer_env / "reports"
    main_mod._REPORT_OUTPUT_DIR = report_dir  # type: ignore[attr-defined]
    main_mod.OUTPUT_DIR = deck_writer_env / "decks"  # type: ignore[attr-defined]
    main_mod._XLSX_OUTPUT_DIR = report_dir  # type: ignore[attr-defined]

    responses = _mock_http_responses()
    mock_client_instance = _make_async_client_mock(responses)

    # Patch httpx.AsyncClient in the main module's namespace so that
    # the lifespan's ``httpx.AsyncClient(...)`` call returns our mock.
    with patch.object(
        main_mod.httpx, "AsyncClient",  # type: ignore[attr-defined]
        return_value=mock_client_instance,
    ):
        with TestClient(main_mod.app) as tc:
            yield tc


# ---------------------------------------------------------------------------
# Test 1: /report with ImageEmbed
# ---------------------------------------------------------------------------

@pytest.mark.integration
def test_report_with_image_embed(client, tmp_path: Path):
    """POST /report with an ImageEmbed should return a .docx with >= 1 image."""
    from docx import Document
    from docx.oxml.ns import qn

    # Create a real PNG file the endpoint can read.
    png_data = _make_minimal_png()
    png_path = tmp_path / "test_image.png"
    png_path.write_bytes(png_data)

    payload = {
        "brief": "Quarterly capital review",
        "dept_id": "cac",
        "title": "Capital Review Q1",
        "images": [
            {
                "path": str(png_path),
                "section_hint": "Background",
                "caption": "Figure 1: Capital overview",
                "width_inches": 4.0,
            }
        ],
    }

    resp = client.post("/report", json=payload)
    assert resp.status_code == 200, resp.text

    data = resp.json()
    assert "file_path" in data
    docx_path = Path(data["file_path"])
    assert docx_path.is_file(), f"Expected .docx at {docx_path}"
    assert docx_path.suffix == ".docx"

    # Verify the docx contains at least one inline image (blip element in XML).
    doc = Document(str(docx_path))
    body_xml = doc.element.body.xml
    blip_tag = qn("a:blip")
    has_image = f"<{blip_tag}" in body_xml or "blip" in body_xml.lower()
    assert has_image, (
        "Expected the .docx to contain at least one embedded image "
        f"(checked for a:blip in body XML). Body XML snippet: {body_xml[:500]!r}"
    )


# ---------------------------------------------------------------------------
# Test 2: /compose with ChartEmbed
# ---------------------------------------------------------------------------

@pytest.mark.integration
def test_compose_with_chart_embed(client):
    """POST /compose with a ChartEmbed should return a .pptx with a picture shape."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    payload = {
        "brief": "Engine performance summary",
        "dept_id": "ic",
        "title": "Engine Performance",
        "charts": [
            {
                "spec": {
                    "kind": "bar",
                    "title": "Return YTD",
                    "series": [{"name": "Return", "data": [12.4, 8.1, 15.2]}],
                    "x_labels": ["VCC", "DAT", "ALT"],
                    "palette": "brooker",
                },
                "slide_index": 0,
                "caption": "Source: ALCO Tracker",
                "width_inches": 5.0,
            }
        ],
    }

    resp = client.post("/compose", json=payload)
    assert resp.status_code == 200, resp.text

    data = resp.json()
    assert "file_path" in data
    pptx_path = Path(data["file_path"])
    assert pptx_path.is_file(), f"Expected .pptx at {pptx_path}"
    assert pptx_path.suffix == ".pptx"

    # Verify the pptx has at least one picture shape across all slides.
    prs = Presentation(str(pptx_path))
    picture_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            # MSO_SHAPE_TYPE.PICTURE == 13
            if shape.shape_type == 13:
                picture_count += 1

    assert picture_count >= 1, (
        f"Expected at least 1 picture shape across all slides, found {picture_count}. "
        f"Slides: {len(prs.slides)}"
    )


# ---------------------------------------------------------------------------
# Test 3: /compose-xlsx with 2 sheets and a formula
# ---------------------------------------------------------------------------

@pytest.mark.integration
def test_compose_xlsx_simple(client):
    """POST /compose-xlsx with 2 sheets + 1 formula; verify both sheets and formula present."""
    import openpyxl

    payload = {
        "spec": {
            "sheets": [
                {
                    "name": "Summary",
                    "headers": ["Item", "Value"],
                    "rows": [
                        ["Revenue", "1000000"],
                        ["Costs", "750000"],
                        ["Profit", ""],
                    ],
                    "formulas": {"B4": "=B2-B3"},
                },
                {
                    "name": "Detail",
                    "headers": ["Category", "Amount", "Pct"],
                    "rows": [
                        ["Fixed", "400000", "53%"],
                        ["Variable", "350000", "47%"],
                    ],
                },
            ],
        },
        "filename": "test_report.xlsx",
        "caller_dept": "cac",
    }

    resp = client.post("/compose-xlsx", json=payload)
    assert resp.status_code == 200, resp.text
    assert resp.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )

    # Parse the returned xlsx bytes directly.
    wb = openpyxl.load_workbook(io.BytesIO(resp.content))

    # Both sheets must be present.
    assert "Summary" in wb.sheetnames, f"'Summary' sheet missing; sheets: {wb.sheetnames}"
    assert "Detail" in wb.sheetnames, f"'Detail' sheet missing; sheets: {wb.sheetnames}"

    # Formula must be in the correct cell.
    ws = wb["Summary"]
    formula_cell = ws["B4"]
    assert formula_cell.value is not None, "Expected a value/formula in Summary!B4"
    assert str(formula_cell.value).startswith("="), (
        f"Expected a formula (starts with '=') in Summary!B4, got: {formula_cell.value!r}"
    )
