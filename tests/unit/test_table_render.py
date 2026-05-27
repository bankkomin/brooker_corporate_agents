"""Unit tests for services/shared/table_render.py and drafter_table_prompt.py.

Run:
    python -m pytest tests/unit/test_table_render.py -v
"""
from __future__ import annotations

import io
import logging

import pytest


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

def _make_simple_spec(**overrides) -> dict:
    base = {
        "headers": ["Name", "Value", "Status"],
        "rows": [
            ["Alpha", "100", "OK"],
            ["Beta", "200", "Watch"],
        ],
    }
    base.update(overrides)
    return base


# ---------------------------------------------------------------------------
# validate_table_spec
# ---------------------------------------------------------------------------

class TestValidateTableSpec:
    def test_validate_rejects_row_length_mismatch(self):
        """Row with 2 cells against 3-column headers must raise ValueError."""
        from services.shared.table_render import validate_table_spec

        spec = {
            "headers": ["A", "B", "C"],
            "rows": [["x", "y"]],  # only 2 cells — mismatch
        }
        with pytest.raises(ValueError, match="Row 0"):
            validate_table_spec(spec)

    def test_validate_coerces_non_str_cells(self):
        """Numbers, booleans, and None should be coerced to str without error."""
        from services.shared.table_render import validate_table_spec

        spec = {
            "headers": ["Label", "Amount", "Active"],
            "rows": [
                ["Engine A", 12.4, True],
                ["Engine B", 0, False],
            ],
        }
        result = validate_table_spec(spec)
        assert result["rows"][0] == ["Engine A", "12.4", "True"]
        assert result["rows"][1] == ["Engine B", "0", "False"]

    def test_validate_rejects_missing_headers(self):
        with pytest.raises(ValueError, match="headers"):
            from services.shared.table_render import validate_table_spec
            validate_table_spec({"rows": [["a"]]})

    def test_validate_rejects_invalid_style(self):
        from services.shared.table_render import validate_table_spec

        with pytest.raises(ValueError, match="style"):
            validate_table_spec({
                "headers": ["A"],
                "rows": [],
                "style": "fancy",
            })

    def test_validate_rejects_column_widths_length_mismatch(self):
        from services.shared.table_render import validate_table_spec

        with pytest.raises(ValueError, match="column_widths_in"):
            validate_table_spec({
                "headers": ["A", "B"],
                "rows": [],
                "column_widths_in": [1.0],  # only 1, but 2 headers
            })

    def test_validate_accepts_valid_spec(self):
        from services.shared.table_render import validate_table_spec

        result = validate_table_spec(_make_simple_spec())
        assert result["headers"] == ["Name", "Value", "Status"]
        assert len(result["rows"]) == 2

    def test_validate_coerces_header_values(self):
        """Headers that are not strings should be coerced."""
        from services.shared.table_render import validate_table_spec

        result = validate_table_spec({"headers": [1, 2, 3], "rows": [[4, 5, 6]]})
        assert result["headers"] == ["1", "2", "3"]
        assert result["rows"][0] == ["4", "5", "6"]


# ---------------------------------------------------------------------------
# add_table_to_docx
# ---------------------------------------------------------------------------

class TestAddTableToDocx:
    def _round_trip(self, doc):
        """Save to bytes and re-open — verifies the XML is valid."""
        from docx import Document

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return Document(buf)

    def test_add_table_to_docx_appends_when_no_hint(self):
        """With no section_hint, a blank doc gets one table appended."""
        from docx import Document

        from services.shared.table_render import add_table_to_docx

        doc = Document()
        result = add_table_to_docx(doc, _make_simple_spec())
        assert result == 1

        reloaded = self._round_trip(doc)
        assert len(reloaded.tables) == 1

    def test_add_table_to_docx_inserts_after_heading(self):
        """Table should appear immediately after the target heading paragraph."""
        from docx import Document

        from services.shared.table_render import add_table_to_docx

        doc = Document()
        doc.add_heading("Introduction", level=1)
        intro_para = doc.add_paragraph("Some intro text.")
        doc.add_heading("Findings", level=1)
        findings_para = doc.add_paragraph("Some findings text.")
        doc.add_heading("Conclusion", level=1)

        add_table_to_docx(doc, _make_simple_spec(), section_hint="Findings")

        # After save/reload verify structure: the table element should follow
        # the "Findings" heading
        reloaded = self._round_trip(doc)
        assert len(reloaded.tables) == 1

        # Walk body children to confirm ordering
        body = reloaded.element.body
        children = list(body)
        ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

        # Find the Findings heading paragraph position
        findings_idx = None
        table_idx = None
        for i, child in enumerate(children):
            if child.tag == f"{{{ns}}}p":
                text = "".join(
                    t.text or "" for t in child.iter(f"{{{ns}}}t")
                )
                if "Findings" in text:
                    style_el = child.find(
                        f".//{{{ns}}}pStyle"
                    )
                    # Confirm it's a heading (style name contains "Heading")
                    if style_el is not None:
                        findings_idx = i
            elif child.tag == f"{{{ns}}}tbl":
                table_idx = i

        assert findings_idx is not None, "Findings heading not found in body"
        assert table_idx is not None, "Table not found in body"
        assert table_idx > findings_idx, (
            f"Table (pos {table_idx}) should come after Findings heading "
            f"(pos {findings_idx})"
        )

        # Additionally: table should come BEFORE the Conclusion heading
        conclusion_idx = None
        for i, child in enumerate(children):
            if child.tag == f"{{{ns}}}p":
                text = "".join(t.text or "" for t in child.iter(f"{{{ns}}}t"))
                if "Conclusion" in text:
                    style_el = child.find(f".//{{{ns}}}pStyle")
                    if style_el is not None:
                        conclusion_idx = i
                        break

        if conclusion_idx is not None:
            assert table_idx < conclusion_idx, (
                f"Table (pos {table_idx}) should come before Conclusion heading "
                f"(pos {conclusion_idx})"
            )

    def test_add_table_to_docx_header_row_bold(self):
        """Header row cells must contain at least one bold run after save/reload."""
        from docx import Document

        from services.shared.table_render import add_table_to_docx

        doc = Document()
        add_table_to_docx(doc, _make_simple_spec())

        reloaded = self._round_trip(doc)
        tbl = reloaded.tables[0]
        hdr_row = tbl.rows[0]

        for cell in hdr_row.cells:
            has_bold = any(
                run.bold
                for para in cell.paragraphs
                for run in para.runs
            )
            assert has_bold, (
                f"Header cell '{cell.text}' has no bold run after save/reload"
            )

    def test_add_table_to_docx_raises_on_invalid_spec(self):
        """ValueError propagates from validate_table_spec."""
        from docx import Document

        from services.shared.table_render import add_table_to_docx

        doc = Document()
        with pytest.raises(ValueError):
            add_table_to_docx(doc, {"headers": ["A", "B"], "rows": [["only_one"]]})

    def test_add_table_to_docx_no_hint_falls_back_gracefully(self):
        """section_hint that matches nothing should still produce a table."""
        from docx import Document

        from services.shared.table_render import add_table_to_docx

        doc = Document()
        doc.add_paragraph("No headings here.")
        result = add_table_to_docx(doc, _make_simple_spec(), section_hint="Missing")
        assert result == 1
        assert len(doc.tables) == 1


# ---------------------------------------------------------------------------
# add_table_to_pptx
# ---------------------------------------------------------------------------

class TestAddTableToPptx:
    def _round_trip(self, prs):
        """Save to bytes and re-open."""
        from pptx import Presentation

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return Presentation(buf)

    def _make_two_slide_prs(self):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        layout = prs.slide_layouts[5]  # blank layout
        prs.slides.add_slide(layout)
        prs.slides.add_slide(layout)
        return prs

    def test_add_table_to_pptx_uses_slide_index(self):
        """Table should be placed on the slide at the specified index (slide 0)."""
        from services.shared.table_render import add_table_to_pptx

        prs = self._make_two_slide_prs()
        result = add_table_to_pptx(prs, _make_simple_spec(), slide_index=0)
        assert result == 1

        reloaded = self._round_trip(prs)
        slide_0 = reloaded.slides[0]
        table_shapes = [s for s in slide_0.shapes if s.has_table]
        assert len(table_shapes) == 1, "Expected exactly 1 table on slide 0"

        # Slide 1 should have no tables
        slide_1 = reloaded.slides[1]
        table_shapes_1 = [s for s in slide_1.shapes if s.has_table]
        assert len(table_shapes_1) == 0, "Slide 1 should have no tables"

    def test_add_table_to_pptx_new_slide_when_no_match(self):
        """Without a matching index or hint, a new slide must be created."""
        from pptx import Presentation

        from services.shared.table_render import add_table_to_pptx

        prs = Presentation()
        original_count = len(prs.slides)

        result = add_table_to_pptx(prs, _make_simple_spec())  # no index, no hint
        assert result == 1

        assert len(prs.slides) == original_count + 1

        # The new (last) slide must have a table
        reloaded = self._round_trip(prs)
        last_slide = reloaded.slides[-1]
        table_shapes = [s for s in last_slide.shapes if s.has_table]
        assert len(table_shapes) == 1

    def test_add_table_to_pptx_header_row_navy(self):
        """Header row should have the Brooker navy fill (#0F3D5C)."""
        from pptx import Presentation
        from pptx.dml.color import RGBColor

        from services.shared.table_render import add_table_to_pptx

        prs = Presentation()
        add_table_to_pptx(prs, _make_simple_spec(), slide_index=None)

        reloaded = self._round_trip(prs)
        last_slide = reloaded.slides[-1]
        table_shapes = [s for s in last_slide.shapes if s.has_table]
        assert table_shapes, "No table found"

        tbl = table_shapes[0].table
        header_cell = tbl.cell(0, 0)
        fill_color = header_cell.fill.fore_color.rgb
        assert fill_color == RGBColor(0x0F, 0x3D, 0x5C), (
            f"Expected navy #0F3D5C, got {fill_color}"
        )

    def test_add_table_to_pptx_out_of_range_index_creates_new_slide(self):
        """slide_index beyond range should create a new slide (with warning)."""
        from pptx import Presentation

        from services.shared.table_render import add_table_to_pptx

        prs = Presentation()
        n = len(prs.slides)
        result = add_table_to_pptx(prs, _make_simple_spec(), slide_index=99)
        assert result == 1
        assert len(prs.slides) == n + 1


# ---------------------------------------------------------------------------
# drafter_table_prompt — extract_tables_from_text
# ---------------------------------------------------------------------------

class TestExtractTablesFromText:
    def test_extract_tables_from_text_finds_one(self):
        """Single valid ```table block is extracted and removed from text."""
        from services.shared.drafter_table_prompt import extract_tables_from_text

        text = (
            "Here is a summary.\n\n"
            "```table\n"
            '{"headers": ["A", "B"], "rows": [["x", "y"]]}\n'
            "```\n\n"
            "End of report."
        )
        cleaned, specs = extract_tables_from_text(text)

        assert len(specs) == 1
        assert specs[0]["headers"] == ["A", "B"]
        assert specs[0]["rows"] == [["x", "y"]]

        # The table block should be gone from cleaned text
        assert "```table" not in cleaned
        assert "```" not in cleaned

        # Prose should remain
        assert "Here is a summary" in cleaned
        assert "End of report" in cleaned

    def test_extract_tables_handles_malformed_json(self):
        """Bad JSON inside a table block: cleaned text removes block, returns []."""
        from services.shared.drafter_table_prompt import extract_tables_from_text

        text = (
            "Before.\n"
            "```table\n"
            "{ this is not json at all }\n"
            "```\n"
            "After."
        )

        with pytest.warns(None) if False else _capture_warnings() as captured:
            cleaned, specs = extract_tables_from_text(text)

        assert specs == [], "Malformed JSON block should yield no specs"
        assert "```table" not in cleaned
        assert "Before." in cleaned
        assert "After." in cleaned

    def test_extract_tables_multiple_blocks(self):
        """Two valid blocks are both extracted."""
        from services.shared.drafter_table_prompt import extract_tables_from_text

        block1 = '```table\n{"headers": ["X"], "rows": [["1"]]}\n```'
        block2 = '```table\n{"headers": ["Y"], "rows": [["2"]]}\n```'
        text = f"Intro.\n\n{block1}\n\nMiddle.\n\n{block2}\n\nEnd."

        cleaned, specs = extract_tables_from_text(text)

        assert len(specs) == 2
        assert specs[0]["headers"] == ["X"]
        assert specs[1]["headers"] == ["Y"]
        assert "```table" not in cleaned

    def test_extract_tables_no_blocks(self):
        """Text with no table blocks returns text unchanged and empty list."""
        from services.shared.drafter_table_prompt import extract_tables_from_text

        text = "Just plain text with no tables."
        cleaned, specs = extract_tables_from_text(text)

        assert specs == []
        assert cleaned == text

    def test_extract_tables_mixed_valid_and_malformed(self):
        """One valid + one malformed: only the valid one is returned."""
        from services.shared.drafter_table_prompt import extract_tables_from_text

        good = '```table\n{"headers": ["A"], "rows": [["v"]]}\n```'
        bad = "```table\nnot json\n```"
        text = f"Start.\n{good}\nMid.\n{bad}\nEnd."

        cleaned, specs = extract_tables_from_text(text)

        assert len(specs) == 1
        assert specs[0]["headers"] == ["A"]
        assert "```table" not in cleaned


# ---------------------------------------------------------------------------
# drafter_table_prompt — table_emission_prompt_snippet
# ---------------------------------------------------------------------------

class TestTableEmissionPromptSnippet:
    def test_table_emission_prompt_snippet_includes_example(self):
        """Snippet must contain a ```table block and the word 'headers'."""
        from services.shared.drafter_table_prompt import table_emission_prompt_snippet

        snippet = table_emission_prompt_snippet()

        assert isinstance(snippet, str)
        assert "```table" in snippet, "Snippet must contain a ```table example"
        assert "headers" in snippet, "Snippet must reference 'headers'"

    def test_snippet_contains_row_example(self):
        """Example in the snippet must have a 'rows' key."""
        from services.shared.drafter_table_prompt import table_emission_prompt_snippet

        snippet = table_emission_prompt_snippet()
        assert '"rows"' in snippet or "'rows'" in snippet

    def test_snippet_mentions_sparing_use(self):
        """Snippet should advise sparse usage (<=3 per doc)."""
        from services.shared.drafter_table_prompt import table_emission_prompt_snippet

        snippet = table_emission_prompt_snippet()
        assert "3" in snippet  # "no more than 3 per document"


# ---------------------------------------------------------------------------
# Context manager helper for warning capture (avoids pytest.warns(None) issue)
# ---------------------------------------------------------------------------

from contextlib import contextmanager


@contextmanager
def _capture_warnings():
    """Minimal context manager that captures logging warnings."""
    records: list[logging.LogRecord] = []

    class _Handler(logging.Handler):
        def emit(self, record: logging.LogRecord) -> None:
            records.append(record)

    handler = _Handler()
    handler.setLevel(logging.WARNING)
    root = logging.getLogger("services.shared.drafter_table_prompt")
    root.addHandler(handler)
    try:
        yield records
    finally:
        root.removeHandler(handler)
