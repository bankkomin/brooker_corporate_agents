"""Unit tests for services/shared/image_embed.py.

Tests are isolated — no network, no MinIO, no vLLM.
All external I/O is mocked via pytest monkeypatch / unittest.mock.

Run with:
    python -m pytest tests/unit/test_image_embed.py -v
"""
from __future__ import annotations

import io
import struct
import zlib
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from docx import Document
from pptx import Presentation


# ---------------------------------------------------------------------------
# Helpers — minimal valid image bytes
# ---------------------------------------------------------------------------

def _make_png(width: int = 1, height: int = 1, r: int = 100, g: int = 150, b: int = 200) -> bytes:
    """Return a minimal, structurally valid PNG for *width* x *height* pixels.

    Creates a true RGB PNG so python-docx / python-pptx can embed it without
    needing a full imaging library at test time.
    """
    sig = b"\x89PNG\r\n\x1a\n"

    # IHDR: width, height, bit_depth=8, color_type=2 (RGB), compression=0,
    #       filter=0, interlace=0
    ihdr_data = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF
    ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc)

    # IDAT: one row of pixels (filter byte 0 + RGB per pixel)
    raw_rows = (b"\x00" + bytes([r, g, b]) * width) * height
    compressed = zlib.compress(raw_rows)
    idat_crc = zlib.crc32(b"IDAT" + compressed) & 0xFFFFFFFF
    idat = struct.pack(">I", len(compressed)) + b"IDAT" + compressed + struct.pack(">I", idat_crc)

    # IEND
    iend_crc = zlib.crc32(b"IEND") & 0xFFFFFFFF
    iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc)

    return sig + ihdr + idat + iend


_SMALL_PNG: bytes = _make_png()          # 69-byte valid PNG
_BAD_BYTES: bytes = b"not an image"      # clearly not PNG or JPEG


# ---------------------------------------------------------------------------
# Import under test (late so the fixture above is defined first)
# ---------------------------------------------------------------------------

from services.shared.image_embed import (  # noqa: E402
    embed_image_in_docx,
    embed_image_in_pptx,
    fetch_image_bytes,
)


# ===========================================================================
# fetch_image_bytes tests
# ===========================================================================

class TestFetchImageBytesFromPath:
    """Local filesystem source."""

    def test_returns_exact_bytes(self, tmp_path: Path) -> None:
        """fetch_image_bytes({"path": ...}) returns the same bytes as the file."""
        img_file = tmp_path / "sample.png"
        img_file.write_bytes(_SMALL_PNG)

        result = fetch_image_bytes({"path": str(img_file)})

        assert result == _SMALL_PNG

    def test_missing_file_raises_value_error(self, tmp_path: Path) -> None:
        with pytest.raises(ValueError, match="not found"):
            fetch_image_bytes({"path": str(tmp_path / "nonexistent.png")})

    def test_bad_magic_raises_value_error(self, tmp_path: Path) -> None:
        """A file that doesn't start with PNG/JPEG magic bytes is rejected."""
        bad_file = tmp_path / "garbage.png"
        bad_file.write_bytes(_BAD_BYTES)

        with pytest.raises(ValueError, match="Unsupported image format"):
            fetch_image_bytes({"path": str(bad_file)})


class TestFetchImageBytesFromUrl:
    """HTTP URL source — network calls are mocked."""

    def _make_streaming_mock(self, body: bytes, content_length: int | None = None) -> MagicMock:
        """Return a mock httpx.Client whose stream() context manager yields *body*."""
        mock_resp = MagicMock()
        mock_resp.headers = {}
        if content_length is not None:
            mock_resp.headers["content-length"] = str(content_length)
        mock_resp.raise_for_status = MagicMock()

        # iter_bytes yields the body as a single chunk
        mock_resp.iter_bytes = MagicMock(return_value=iter([body]))

        stream_ctx = MagicMock()
        stream_ctx.__enter__ = MagicMock(return_value=mock_resp)
        stream_ctx.__exit__ = MagicMock(return_value=False)

        mock_client = MagicMock()
        # HEAD returns 405 (not allowed) — that path is swallowed gracefully
        mock_client.head.side_effect = Exception("HEAD not supported")
        mock_client.stream = MagicMock(return_value=stream_ctx)

        client_ctx = MagicMock()
        client_ctx.__enter__ = MagicMock(return_value=mock_client)
        client_ctx.__exit__ = MagicMock(return_value=False)
        return client_ctx

    def test_returns_image_bytes(self) -> None:
        client_ctx = self._make_streaming_mock(_SMALL_PNG)
        with patch("httpx.Client", return_value=client_ctx):
            result = fetch_image_bytes({"url": "https://example.com/img.png"})
        assert result == _SMALL_PNG

    def test_rejects_oversized_content_length(self) -> None:
        """When Content-Length exceeds 25 MB, raise before downloading body."""
        oversized = 30 * 1024 * 1024  # 30 MB
        client_ctx = self._make_streaming_mock(_SMALL_PNG, content_length=oversized)
        with patch("httpx.Client", return_value=client_ctx):
            with pytest.raises(ValueError, match="too large"):
                fetch_image_bytes({"url": "https://example.com/big.png"})

    def test_invalid_scheme_raises(self) -> None:
        with pytest.raises(ValueError, match="http"):
            fetch_image_bytes({"url": "ftp://example.com/img.png"})


class TestFetchImageBytesBadMagic:
    """Magic-byte validation on path source."""

    def test_garbage_file_raises_value_error(self, tmp_path: Path) -> None:
        """A file whose first bytes don't match PNG or JPEG is rejected."""
        garbage = tmp_path / "bad.png"
        garbage.write_bytes(_BAD_BYTES)

        with pytest.raises(ValueError, match="Unsupported image format"):
            fetch_image_bytes({"path": str(garbage)})

    def test_gif_rejected(self, tmp_path: Path) -> None:
        """GIF magic bytes (GIF87a) are not in our allow-list."""
        gif_file = tmp_path / "anim.gif"
        gif_file.write_bytes(b"GIF87a" + b"\x00" * 50)

        with pytest.raises(ValueError, match="Unsupported image format"):
            fetch_image_bytes({"path": str(gif_file)})


class TestFetchImageBytesUnknownKey:
    """Bad ImageSource dicts are rejected immediately."""

    def test_empty_dict_raises(self) -> None:
        with pytest.raises(ValueError, match="non-empty dict"):
            fetch_image_bytes({})

    def test_unknown_key_raises(self) -> None:
        with pytest.raises(ValueError, match="Unrecognised ImageSource"):
            fetch_image_bytes({"s3_key": "bucket/img.png"})  # type: ignore[arg-type]


# ===========================================================================
# embed_image_in_docx tests
# ===========================================================================

class TestEmbedImageInDocxAppendsWhenNoHint:
    """With no section_hint the image should appear as the last shape in the doc."""

    def test_inline_shape_exists_after_save_reload(self, tmp_path: Path) -> None:
        doc = Document()
        doc.add_paragraph("Some preamble text.")

        result = embed_image_in_docx(doc, {"path": _png_path(tmp_path)})

        assert result == 1

        # Save and reload to verify the relationship is properly serialised.
        out = tmp_path / "out.docx"
        doc.save(str(out))
        reloaded = Document(str(out))

        # Count inline shapes across all paragraphs.
        shapes = _count_inline_shapes(reloaded)
        assert shapes == 1, f"Expected 1 inline shape, found {shapes}"

    def test_returns_one_on_success(self, tmp_path: Path) -> None:
        doc = Document()
        assert embed_image_in_docx(doc, {"path": _png_path(tmp_path)}) == 1


class TestEmbedImageInDocxInsertsAfterHeading:
    """section_hint should place the image right after the matched heading."""

    def test_image_lands_at_correct_paragraph_index(self, tmp_path: Path) -> None:
        doc = Document()
        doc.add_heading("Introduction", level=1)
        doc.add_paragraph("Intro body text.")
        doc.add_heading("Executive Summary", level=1)
        doc.add_paragraph("Exec summary body text.")
        doc.add_heading("Conclusion", level=1)

        embed_image_in_docx(
            doc,
            {"path": _png_path(tmp_path)},
            section_hint="Executive Summary",
        )

        # Find paragraph index of "Executive Summary" heading.
        paragraphs = doc.paragraphs
        heading_idx = next(
            i for i, p in enumerate(paragraphs)
            if "Executive Summary" in p.text
        )

        # The paragraph immediately after the heading should contain our image
        # (an inline shape means the run's XML has a <w:drawing> element).
        next_para = paragraphs[heading_idx + 1]
        assert _para_has_inline_shape(next_para), (
            f"Expected an inline shape in paragraph {heading_idx + 1!r}, "
            f"got text: {next_para.text!r}"
        )

    def test_no_match_falls_back_to_append(self, tmp_path: Path) -> None:
        """When section_hint matches nothing, image is appended — no error."""
        doc = Document()
        doc.add_heading("Only Heading", level=1)

        result = embed_image_in_docx(
            doc,
            {"path": _png_path(tmp_path)},
            section_hint="Nonexistent Section XYZ",
        )
        assert result == 1
        shapes = _count_inline_shapes(doc)
        assert shapes == 1


class TestEmbedImageInDocxWithCaption:
    """Caption paragraph should appear right after the image and be italic."""

    def test_caption_text_is_present(self, tmp_path: Path) -> None:
        doc = Document()
        caption_text = "Figure 1: Corporate overview chart"

        embed_image_in_docx(
            doc,
            {"path": _png_path(tmp_path)},
            caption=caption_text,
        )

        # The caption paragraph should exist somewhere in the document.
        all_text = " ".join(p.text for p in doc.paragraphs)
        assert caption_text in all_text, f"Caption not found in: {all_text!r}"

    def test_caption_paragraph_is_italic(self, tmp_path: Path) -> None:
        doc = Document()
        caption_text = "Quarterly performance chart"

        embed_image_in_docx(
            doc,
            {"path": _png_path(tmp_path)},
            caption=caption_text,
        )

        # Find the caption paragraph and check its run's italic flag.
        caption_para = next(
            p for p in doc.paragraphs if caption_text in p.text
        )
        assert caption_para.runs, "Caption paragraph has no runs"
        assert caption_para.runs[0].italic, (
            "First run of caption paragraph is not italic"
        )


# ===========================================================================
# embed_image_in_pptx tests
# ===========================================================================

class TestEmbedImageInPptxUsesSlideIndex:
    """slide_index should target the specified slide (0-based)."""

    def test_picture_shape_on_correct_slide(self, tmp_path: Path) -> None:
        prs = Presentation()
        # Add two blank slides.
        layout = prs.slide_layouts[5]
        prs.slides.add_slide(layout)
        prs.slides.add_slide(layout)
        assert len(prs.slides) == 2

        result = embed_image_in_pptx(
            prs,
            {"path": _png_path(tmp_path)},
            slide_index=1,
        )

        assert result == 1
        # Slide at index 1 should now have a picture shape; slide 0 should not.
        slide_1_pics = _count_picture_shapes(prs.slides[1])
        slide_0_pics = _count_picture_shapes(prs.slides[0])
        assert slide_1_pics == 1, f"Expected 1 picture on slide 1, found {slide_1_pics}"
        assert slide_0_pics == 0, f"Expected 0 pictures on slide 0, found {slide_0_pics}"

    def test_slide_index_zero(self, tmp_path: Path) -> None:
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])

        embed_image_in_pptx(prs, {"path": _png_path(tmp_path)}, slide_index=0)

        assert _count_picture_shapes(prs.slides[0]) == 1


class TestEmbedImageInPptxNewSlideWhenNoMatch:
    """With neither slide_index nor slide_title_hint, a new slide is appended."""

    def test_new_slide_created_with_picture(self, tmp_path: Path) -> None:
        prs = Presentation()
        initial_count = len(prs.slides)  # 0

        result = embed_image_in_pptx(
            prs,
            {"path": _png_path(tmp_path)},
            # No slide_index, no slide_title_hint
        )

        assert result == 1
        assert len(prs.slides) == initial_count + 1, (
            f"Expected a new slide to be added (was {initial_count}, "
            f"now {len(prs.slides)})"
        )
        new_slide = prs.slides[initial_count]
        assert _count_picture_shapes(new_slide) == 1

    def test_out_of_range_index_creates_new_slide(self, tmp_path: Path) -> None:
        """An out-of-range slide_index should fall through to appending."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])
        initial_count = len(prs.slides)

        embed_image_in_pptx(prs, {"path": _png_path(tmp_path)}, slide_index=99)

        assert len(prs.slides) == initial_count + 1

    def test_title_hint_match_uses_existing_slide(self, tmp_path: Path) -> None:
        """slide_title_hint matches an existing slide — no new slide created."""
        from pptx.util import Pt

        prs = Presentation()
        layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(layout)
        # Set title text on the slide.
        if slide.shapes.title:
            slide.shapes.title.text = "Financials Overview"

        initial_count = len(prs.slides)

        embed_image_in_pptx(
            prs,
            {"path": _png_path(tmp_path)},
            slide_title_hint="Financials",  # substring match
        )

        assert len(prs.slides) == initial_count, (
            "No new slide should be created when title_hint matches"
        )
        assert _count_picture_shapes(prs.slides[0]) == 1


# ===========================================================================
# Private helpers for tests
# ===========================================================================

def _png_path(tmp_path: Path) -> str:
    """Write the canonical test PNG to *tmp_path* and return its path string."""
    p = tmp_path / "test.png"
    p.write_bytes(_SMALL_PNG)
    return str(p)


def _count_inline_shapes(doc) -> int:
    """Count paragraphs that contain at least one inline image (<w:drawing>)."""
    from lxml import etree

    # python-docx doesn't expose inline shapes on Document directly;
    # we check each paragraph's XML for the <w:drawing> element.
    ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    drawing_tag = f"{{{ns}}}drawing"

    count = 0
    for para in doc.paragraphs:
        for elem in para._p.iter():
            if elem.tag == drawing_tag:
                count += 1
                break  # one match per paragraph is enough
    return count


def _para_has_inline_shape(para) -> bool:
    """Return True if the paragraph's XML contains a <w:drawing> element."""
    ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    drawing_tag = f"{{{ns}}}drawing"
    return any(elem.tag == drawing_tag for elem in para._p.iter())


def _count_picture_shapes(slide) -> int:
    """Count PICTURE shapes on a pptx slide."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return sum(
        1 for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
